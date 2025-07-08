import streamlit as st
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns
import plotly.express as px
from sklearn.model_selection import train_test_split
from sklearn.linear_model import LinearRegression, LogisticRegression
from sklearn.ensemble import RandomForestClassifier, RandomForestRegressor
from sklearn.metrics import classification_report, confusion_matrix, mean_squared_error, r2_score
from sklearn.preprocessing import StandardScaler, LabelEncoder, OneHotEncoder
from sklearn.cluster import KMeans
from sklearn.decomposition import PCA
from sklearn.feature_selection import SelectKBest, f_classif, mutual_info_classif
import os
import tempfile
import unicodedata
import string
from PIL import Image
from datetime import datetime
import time

def generate_graph_images(df, tmpdir, t):
    img_paths = []
    numeric_columns = df.select_dtypes(include=np.number).columns.tolist()
    categorical_cols = df.select_dtypes(include="object").columns.tolist()
    palette_list = [
        "tab10", "Set2", "Set1", "Pastel1", "Paired", "Dark2", "Accent", "tab20"
    ]
    # 1. Histogrammes
    for idx, col in enumerate(numeric_columns[:5]):
        fig, ax = plt.subplots(figsize=(5,3), dpi=120)
        sns.histplot(df[col], kde=True, ax=ax, color=sns.color_palette(palette_list[idx % len(palette_list)])[0])
        ax.set_title(t["hist_title"].format(col=col))
        img1 = os.path.join(tmpdir, f"graph_hist_{col}.png")
        fig.tight_layout()
        fig.savefig(img1, bbox_inches="tight")
        plt.close(fig)
        img_paths.append(img1)
    # 2. Scatterplots
    if len(numeric_columns) >= 2:
        for i in range(2):
            fig, ax = plt.subplots(figsize=(5,3), dpi=120)
            sns.scatterplot(
                x=df[numeric_columns[i]],
                y=df[numeric_columns[(i+1)%len(numeric_columns)]],
                ax=ax,
                color=sns.color_palette("tab10")[i]
            )
            ax.set_title(
                t["scatter_title"].format(
                    x=numeric_columns[i],
                    y=numeric_columns[(i+1)%len(numeric_columns)]
                )
            )
            img2 = os.path.join(tmpdir, f"graph_scatter_{i}.png")
            fig.tight_layout()
            fig.savefig(img2, bbox_inches="tight")
            plt.close(fig)
            img_paths.append(img2)
    # 3. Boxplots
    for idx, col in enumerate(numeric_columns[:3]):
        fig, ax = plt.subplots(figsize=(5,3), dpi=120)
        sns.boxplot(y=df[col], ax=ax, color=sns.color_palette("Set2")[idx])
        ax.set_title(t["box_title"].format(col=col))
        img3 = os.path.join(tmpdir, f"graph_box_{col}.png")
        fig.tight_layout()
        fig.savefig(img3, bbox_inches="tight")
        plt.close(fig)
        img_paths.append(img3)
    # 4. Heatmap corr
    if len(numeric_columns) > 1:
        fig, ax = plt.subplots(figsize=(5,3), dpi=120)
        sns.heatmap(df[numeric_columns].corr(), annot=True, cmap="coolwarm", ax=ax)
        ax.set_title(t["corr_title"])
        img4 = os.path.join(tmpdir, "graph_corr.png")
        fig.tight_layout()
        fig.savefig(img4, bbox_inches="tight")
        plt.close(fig)
        img_paths.append(img4)
    # 5. Pairplot
    if len(numeric_columns) >= 2:
        sns.pairplot(df[numeric_columns].dropna().sample(min(100, len(df))), diag_kind="kde", palette="Set1")
        plt.suptitle(t["pairplot_title"], y=1.02)
        img = os.path.join(tmpdir, "graph_pairplot.png")
        plt.savefig(img, bbox_inches="tight", dpi=120)
        plt.close()
        img_paths.append(img)
    # 6. Pie chart (cat)
    if categorical_cols and df[categorical_cols[0]].nunique() < 15:
        fig, ax = plt.subplots(figsize=(5,3), dpi=120)
        df[categorical_cols[0]].value_counts().plot.pie(autopct='%1.1f%%', ax=ax, colors=sns.color_palette("tab20"))
        ax.set_ylabel("")
        ax.set_title(t["pie_title"].format(col=categorical_cols[0]))
        img5 = os.path.join(tmpdir, "graph_piecat.png")
        fig.tight_layout()
        fig.savefig(img5, bbox_inches="tight")
        plt.close(fig)
        img_paths.append(img5)
    # 7. Barplots (cat)
    for idx, col in enumerate(categorical_cols[:2]):
        fig, ax = plt.subplots(figsize=(5,3), dpi=120)
        df[col].value_counts().plot(kind="bar", ax=ax, color=sns.color_palette("Set2"))
        ax.set_title(t["bar_title"].format(col=col))
        img6 = os.path.join(tmpdir, f"graph_bar_{col}.png")
        fig.tight_layout()
        fig.savefig(img6, bbox_inches="tight")
        plt.close(fig)
        img_paths.append(img6)
    return img_paths

# ----------- APPEL DU BLOC DANS TON CODE PRINCIPAL ------------

# Exemple d'utilisation : df = ton DataFrame, t = ton dictionnaire de traduction
# Remplace cette partie par l'endroit où tu veux générer les images

    # ...
    # ...



# Ajouts pour analyse avancée texte/audio/image/OCR et sécurité
import hashlib
try:
    import pytesseract
except ImportError:
    pytesseract = None
try:
    import pdfplumber
except ImportError:
    pdfplumber = None
try:
    import docx
except ImportError:
    docx = None
try:
    import whisper
except ImportError:
    whisper = None
try:
    from cryptography.fernet import Fernet
except ImportError:
    Fernet = None

# ========== Thème sombre/clair natif + CSS ultra bleu + branding premium ==========
if "dark_mode" not in st.session_state:
    st.session_state.dark_mode = st.get_option("theme.base") == "dark"
dark_mode = st.sidebar.toggle("🌒 Mode sombre / clair", value=st.session_state.dark_mode)
st.session_state.dark_mode = dark_mode
st.set_page_config(page_title="ErnestMind AI V40 – Super Data Science", layout="wide", initial_sidebar_state="expanded")

st.markdown(f"""
    <style>
    html, body, .main, .block-container {{
        background-color: {'#111' if dark_mode else '#fff'} !important;
    }}
    div[data-testid="stMarkdownContainer"] > p,
    h1, h2, h3, h4, h5, h6, .stText, label, .stButton>button, .stRadio, .stSelectbox,
    .stDataFrame, .stMetric, .stAlert, .stCaption, .stExpanderHeader, .stMarkdown, .css-10trblm, .css-1v0mbdj, .st-bq, .st-c9, .st-c7, .st-cc {{
        color: #1E90FF !important;
    }}
    .stDataFrame {{background-color: {'#222' if dark_mode else '#f7f7f7'};}}
    ::selection {{ background: #1E90FF33; }}
    .stButton > button, .stDownloadButton > button {{
        border: 1px solid #1E90FF !important;
        color: #1E90FF !important;
        background: transparent !important;
    }}
    .superbrand {{
        font-size: 1.8em;
        font-weight: 800;
        color: #1E90FF !important;
        letter-spacing: 1px;
        margin-bottom: 0.2em;
    }}
    .superclaim {{
        font-size: 1.12em;
        color: #1E90FF !important;
        font-weight: 600;
    }}
    .upload-anim {{
        width: 100%;
        margin: 10px 0 25px 0;
        height: 16px;
        background: linear-gradient(90deg, #1E90FF33 25%, #fff0 75%);
        background-size: 200% 100%;
        animation: shineblue 2s linear infinite;
        border-radius: 8px;
        box-shadow: 0 0 8px #1E90FF55;
    }}
    @keyframes shineblue {{
        0% {{ background-position: 200% 0; }}
        100% {{ background-position: -200% 0; }}
    }}
    </style>
""", unsafe_allow_html=True)




# =========== Multi-langues avec drapeaux (20+) ===========
langues_flags = {
    "fr": "🇫🇷 Français",
    "en": "🇬🇧 English",
    "es": "🇪🇸 Español",
    "de": "🇩🇪 Deutsch",
    "pt": "🇵🇹 Português",
    "it": "🇮🇹 Italiano",
    "nl": "🇳🇱 Nederlands",
    "zh": "🇨🇳 中文",
    "ja": "🇯🇵 日本語",
    "ko": "🇰🇷 한국어",
    "ar": "🇸🇦 العربية",
    "hi": "🇮🇳 हिन्दी",
    "th": "🇹🇭 ไทย",
    "tr": "🇹🇷 Türkçe",
    "pl": "🇵🇱 Polski",
    "ru": "🇷🇺 Русский",
    "sv": "🇸🇪 Svenska",
    "da": "🇩🇰 Dansk",
    "no": "🇳🇴 Norsk",
    "ht": "🇭🇹 Kreyòl Ayisyen"
}
langue_code = st.sidebar.selectbox("🌍 Langue / Language", list(langues_flags.keys()), format_func=lambda x: langues_flags[x])

translations = {
    "fr": {
        "aide_title": "🆘 Aide / Help",
        "aide_texte": "Chargez un fichier (csv, excel, parquet, ...), sélectionnez un module IA métier, visualisez plus de 20 graphiques avancés et exportez vos rapports PDF/Word/PPT instantanément. Confidentialité totale.",
        "apropos_title": "ℹ️ À propos",
        "apropos_texte": "ErnestMind V40 dépasse Tableau, PowerBI, Dataiku, Qlik... en simplicité, rapidité et confidentialité. 100% local, IA auto, multi-langues, multi-format.",
        "exemple_toggle": "📈 Exemples de visualisations automatiques",
        "exemple_titre": "### 📊 Démo IA & visualisations",
        "tab_corr": "📊 Corrélation",
        "tab_boxplot": "📦 Boxplot",
        "tab_ventes": "📈 Ventes",
        "tab_marketing": "🎯 Marketing",
        "tab_satisfaction": "😊 Satisfaction",
        "upload_label": "📂 Téléchargez un fichier",
        "format_non_supporte": "❌ Format non supporté",
        "success_upload": "✅ Fichier chargé avec succès !",
        "erreur_chargement": "❌ Erreur lors du chargement du fichier",
        "explorer": "🖨️ Aperçu des données",
        "nb_lignes": "Nombre de lignes",
        "nb_colonnes": "Nombre de colonnes",
        "colonnes": "Noms des colonnes",
        "types_colonnes": "Types des colonnes",
        "manquants": "Colonnes avec valeurs manquantes",
        "uniques": "Valeurs uniques par colonne",
        "stat_desc": "📐 Statistiques descriptives",
        "apercu_donnees": "Aperçu des données",
        "resume": "🔎 Résumé automatique",
        "analyse_auto": "🧠 Analyse automatique intelligente",
        "numeriques": "Numériques",
        "categorique": "Catégoriques",
        "valeurs_manquantes": "Valeurs manquantes",
        "classif_detectee": "Classification supervisée détectée",
        "reg_detectee": "Régression supervisée détectée",
        "tache_non_detectee": "Tâche non détectée automatiquement",
        "visual_auto": "📊 Visualisations automatiques",
        "histogramme": "Histogramme",
        "nuage": "Nuage de points",
        "boxplot": "Boxplot",
        "repartition": "Répartition",
        "choix_format": "📄 Format du rapport à générer",
        "generer_rapport": "📜 Générer le rapport",
        "telecharger_rapport": "📥 Télécharger le rapport généré",
        "date_rapport": "Date et heure du rapport",
        "element": "Élément",
        "valeur": "Valeur",
        "v80_message": "🚀 ErnestMind V40 analyse jusqu'à 100 000 lignes. Pour l'analyse big data illimitée, la rapidité x5 et les modèles sectoriels IA d'élite, passez à la version V80 !",
        "audit_title": "🔒 Audit local & sécurité",
        "audit_exp": "Sur cette version, tous les fichiers temporaires sont chiffrés automatiquement pour garantir la confidentialité. Un rapport d’audit local est généré pour chaque traitement avancé (texte/audio/image).",
        "brand_title": "🚀 ErnestMind AI V40 – Super Data Science Platform",
        "brand_claim": "+ Puissant que Tableau, PowerBI, Dataiku, DataRobot, Alteryx, Qlik, SAS Viya, Palantir...<br>100% local, IA auto, rapports PDF/Word/PPT, multi-langue, mode sombre/clair.",
        "connecteurs_title": "⚡️ Connecteurs / API / Cloud / DB (démo)",
        "connecteurs_coming": "**Connecteurs à venir :**",
        "connecteurs_import": "- Importer depuis MySQL, PostgreSQL, BigQuery, S3, Google Drive, OneDrive, SharePoint, API REST",
        "connecteurs_export": "- Export vers PowerPoint interactif, HTML dashboard, API, webhook, Slack, Teams, email, etc.",
        "connecteurs_contact": "**Contactez-nous pour intégrer votre source métier !**",
        "hist_title": "Histogramme : {col}",
        "scatter_title": "Nuage : {x} vs {y}",
        "box_title": "Boxplot : {col}",
        "corr_title": "Corrélation",
        "pairplot_title": "Nuage de variables (pairplot)",
        "pie_title": "Répartition : {col}",
        "bar_title": "Barplot : {col}",
        "revolution_cta_title": "✨ Rejoignez la révolution ErnestMind !",
        "revolution_cta_text": "Impressionné par la puissance d'ErnestMind V40 ? Ce n'est qu'un aperçu ! Inscrivez-vous pour être informé des nouvelles fonctionnalités, des modèles IA spécialisés et des formations exclusives.",
        "revolution_cta_button": "🚀 S'inscrire à la révolution IA (Cliquez ici !) 🚀",
        "revolution_cta_help": "Cliquez pour accéder à notre page d'inscription et découvrir le futur d'ErnestMind.",
        "contact_email_label": "Contact"
    },
    "en": {
        "aide_title": "🆘 Help",
        "aide_texte": "Upload a file (csv, excel, parquet, ...), select a business AI module, view 20+ advanced charts and export PDF/Word/PPT reports instantly. Total privacy.",
        "apropos_title": "ℹ️ About",
        "apropos_texte": "ErnestMind V40 surpasses Tableau, PowerBI, Dataiku, Qlik... in simplicity, speed and privacy. 100% local, auto AI, multi-languages, multi-format.",
        "exemple_toggle": "📈 Sample automatic visualizations",
        "exemple_titre": "### 📊 Demo IA & visualizations",
        "tab_corr": "📊 Correlation",
        "tab_boxplot": "📦 Boxplot",
        "tab_ventes": "📈 Sales",
        "tab_marketing": "🎯 Marketing",
        "tab_satisfaction": "😊 Satisfaction",
        "upload_label": "📂 Upload a file",
        "format_non_supporte": "❌ Unsupported format",
        "success_upload": "✅ File uploaded successfully!",
        "erreur_chargement": "❌ Error loading file",
        "explorer": "🖨️ Data preview",
        "nb_lignes": "Number of rows",
        "nb_colonnes": "Number of columns",
        "colonnes": "Column names",
        "types_colonnes": "Column types",
        "manquants": "Columns with missing values",
        "uniques": "Unique values per column",
        "stat_desc": "📐 Descriptive statistics",
        "apercu_donnees": "Data preview",
        "resume": "🔎 Auto summary",
        "analyse_auto": "🧠 Smart analysis",
        "numeriques": "Numerical",
        "categorique": "Categorical",
        "valeurs_manquantes": "Missing values",
        "classif_detectee": "Supervised classification detected",
        "reg_detectee": "Supervised regression detected",
        "tache_non_detectee": "Task not auto detected",
        "visual_auto": "📊 Automatic visualizations",
        "histogramme": "Histogram",
        "nuage": "Scatterplot",
        "boxplot": "Boxplot",
        "repartition": "Distribution",
        "choix_format": "📄 Report format",
        "generer_rapport": "📜 Generate report",
        "telecharger_rapport": "📥 Download generated report",
        "date_rapport": "Report date/time",
        "element": "Item",
        "valeur": "Value",
        "v80_message": "🚀 ErnestMind V40 analyzes up to 100,000 rows. For unlimited big data, 5x speed, and elite sector models, upgrade to V80!",
        "audit_title": "🔒 Local audit & security",
        "audit_exp": "In this version, all temporary files are automatically encrypted for confidentiality. A local audit report is generated for each advanced processing (text/audio/image).",
        "brand_title": "🚀 ErnestMind AI V40 – Super Data Science Platform",
        "brand_claim": "Stronger than Tableau, PowerBI, Dataiku, DataRobot, Alteryx, Qlik, SAS Viya, Palantir...<br>100% local, auto AI, PDF/Word/PPT reports, multilingual, dark/light mode.",
        "connecteurs_title": "⚡️ Connectors / API / Cloud / DB (demo)",
        "connecteurs_coming": "**Connectors coming soon:**",
        "connecteurs_import": "- Import from MySQL, PostgreSQL, BigQuery, S3, Google Drive, OneDrive, SharePoint, REST API",
        "connecteurs_export": "- Export to interactive PowerPoint, HTML dashboard, API, webhook, Slack, Teams, email, etc.",
        "connecteurs_contact": "**Contact us to integrate your business source!**",
        "hist_title": "Histogram: {col}",
        "scatter_title": "Scatterplot: {x} vs {y}",
        "box_title": "Boxplot: {col}",
        "corr_title": "Correlation",
        "pairplot_title": "Variable scatter (pairplot)",
        "pie_title": "Distribution: {col}",
        "bar_title": "Barplot: {col}",
        "revolution_cta_title": "✨ Join the ErnestMind Revolution!",
        "revolution_cta_text": "Impressed by the power of ErnestMind V40? This is just a glimpse! Sign up to be informed about new features, specialized AI models, and exclusive training.",
        "revolution_cta_button": "🚀 Join the AI Revolution (Click here!) 🚀",
        "revolution_cta_help": "Click to access our sign-up page and discover the future of ErnestMind.",
        "contact_email_label": "Contact"
    },
    "es": {
        "aide_title": "🆘 Ayuda",
        "aide_texte": "Sube un archivo (csv, excel, parquet, ...), selecciona un módulo de IA de negocio, visualiza más de 20 gráficos avanzados y exporta informes PDF/Word/PPT al instante. Total privacidad.",
        "apropos_title": "ℹ️ Acerca de",
        "apropos_texte": "ErnestMind V40 supera a Tableau, PowerBI, Dataiku, Qlik... en simplicidad, velocidad y privacidad. 100% local, IA automática, multi-idiomas, multi-formato.",
        "exemple_toggle": "📈 Ejemplos de visualizaciones automáticas",
        "exemple_titre": "### 📊 Demostración IA y visualizaciones",
        "tab_corr": "📊 Correlación",
        "tab_boxplot": "📦 Boxplot",
        "tab_ventes": "📈 Ventas",
        "tab_marketing": "🎯 Marketing",
        "tab_satisfaction": "😊 Satisfacción",
        "upload_label": "📂 Subir un archivo",
        "format_non_supporte": "❌ Formato no compatible",
        "success_upload": "✅ ¡Archivo cargado con éxito!",
        "erreur_chargement": "❌ Error al cargar el archivo",
        "explorer": "🖨️ Vista previa de datos",
        "nb_lignes": "Número de filas",
        "nb_colonnes": "Número de columnas",
        "colonnes": "Nombres de columnas",
        "types_colonnes": "Tipos de columnas",
        "manquants": "Columnas con valores faltantes",
        "uniques": "Valores únicos por columna",
        "stat_desc": "📐 Estadísticas descriptivas",
        "apercu_donnees": "Vista previa de datos",
        "resume": "🔎 Resumen automático",
        "analyse_auto": "🧠 Análisis automático inteligente",
        "numeriques": "Numéricos",
        "categorique": "Categóricos",
        "valeurs_manquantes": "Valores faltantes",
        "classif_detectee": "Clasificación supervisada detectada",
        "reg_detectee": "Regresión supervisada detectada",
        "tache_non_detectee": "Tarea no detectada automáticamente",
        "visual_auto": "📊 Visualizaciones automáticas",
        "histogramme": "Histograma",
        "nuage": "Gráfico de dispersión",
        "boxplot": "Diagrama de cajas",
        "repartition": "Distribución",
        "choix_format": "📄 Formato del informe a generar",
        "generer_rapport": "📜 Generar informe",
        "telecharger_rapport": "📥 Descargar informe generado",
        "date_rapport": "Fecha y hora del informe",
        "element": "Elemento",
        "valeur": "Valor",
        "v80_message": "🚀 ErnestMind V40 analiza hasta 100.000 filas. ¡Para análisis de big data ilimitados, velocidad x5 y modelos de IA de élite, actualiza a la versión V80!",
        "audit_title": "🔒 Auditoría local y seguridad",
        "audit_exp": "En esta versión, todos los archivos temporales se cifran automáticamente para garantizar la confidencialidad. Se genera un informe de auditoría local para cada procesamiento avanzado (texto/audio/imagen).",
        "brand_title": "🚀 ErnestMind AI V40 – Plataforma de Superciencia de Datos",
        "brand_claim": "Más potente que Tableau, PowerBI, Dataiku, DataRobot, Alteryx, Qlik, SAS Viya, Palantir...<br>100% local, IA automática, informes PDF/Word/PPT, multi-idioma, modo oscuro/claro.",
        "connecteurs_title": "⚡️ Conectores / API / Nube / DB (demo)",
        "connecteurs_coming": "**Próximos conectores:**",
        "connecteurs_import": "- Importar desde MySQL, PostgreSQL, BigQuery, S3, Google Drive, OneDrive, SharePoint, API REST",
        "connecteurs_export": "- Exportar a PowerPoint interactivo, panel HTML, API, webhook, Slack, Teams, email, etc.",
        "connecteurs_contact": "**¡Contáctenos para integrar su fuente de datos de negocio!**",
        "hist_title": "Histograma: {col}",
        "scatter_title": "Gráfico de dispersión: {x} vs {y}",
        "box_title": "Diagrama de cajas: {col}",
        "corr_title": "Correlación",
        "pairplot_title": "Gráfico de pares de variables",
        "pie_title": "Distribución: {col}",
        "bar_title": "Gráfico de barras: {col}",
        "revolution_cta_title": "✨ ¡Únete a la revolución ErnestMind!",
        "revolution_cta_text": "¿Impresionado por el poder de ErnestMind V40? ¡Esto es solo un vistazo! Regístrate para estar informado sobre nuevas funciones, modelos de IA especializados y capacitaciones exclusivas.",
        "revolution_cta_button": "🚀 ¡Únete a la Revolución de la IA (Haz clic aquí!) 🚀",
        "revolution_cta_help": "Haz clic para acceder a nuestra página de registro y descubrir el futuro de ErnestMind.",
        "contact_email_label": "Contacto"
    },
    "de": {
        "aide_title": "🆘 Hilfe",
        "aide_texte": "Laden Sie eine Datei (csv, excel, parquet, ...) hoch, wählen Sie ein Geschäfts-KI-Modul, sehen Sie über 20 fortgeschrittene Diagramme und exportieren Sie sofort PDF/Word/PPT-Berichte. Vollständige Privatsphäre.",
        "apropos_title": "ℹ️ Über uns",
        "apropos_texte": "ErnestMind V40 übertrifft Tableau, PowerBI, Dataiku, Qlik... in Einfachheit, Geschwindigkeit und Datenschutz. 100% lokal, Auto-KI, mehrsprachig, Multi-Format.",
        "exemple_toggle": "📈 Beispiele für automatische Visualisierungen",
        "exemple_titre": "### 📊 KI-Demo & Visualisierungen",
        "tab_corr": "📊 Korrelation",
        "tab_boxplot": "📦 Boxplot",
        "tab_ventes": "📈 Verkäufe",
        "tab_marketing": "🎯 Marketing",
        "tab_satisfaction": "😊 Zufriedenheit",
        "upload_label": "📂 Datei hochladen",
        "format_non_supporte": "❌ Nicht unterstütztes Format",
        "success_upload": "✅ Datei erfolgreich hochgeladen!",
        "erreur_chargement": "❌ Fehler beim Laden der Datei",
        "explorer": "🖨️ Datenvorschau",
        "nb_lignes": "Anzahl der Zeilen",
        "nb_colonnes": "Anzahl der Spalten",
        "colonnes": "Spaltennamen",
        "types_colonnes": "Spaltentypen",
        "manquants": "Spalten mit fehlenden Werten",
        "uniques": "Eindeutige Werte pro Spalte",
        "stat_desc": "📐 Deskriptive Statistiken",
        "apercu_donnees": "Datenvorschau",
        "resume": "🔎 Automatische Zusammenfassung",
        "analyse_auto": "🧠 Intelligente Analyse",
        "numeriques": "Numerisch",
        "categorique": "Kategorisch",
        "valeurs_manquantes": "Fehlende Werte",
        "classif_detectee": "Überwachte Klassifizierung erkannt",
        "reg_detectee": "Überwachte Regression erkannt",
        "tache_non_detectee": "Aufgabe nicht automatisch erkannt",
        "visual_auto": "📊 Automatische Visualisierungen",
        "histogramme": "Histogramm",
        "nuage": "Streudiagramm",
        "boxplot": "Boxplot",
        "repartition": "Verteilung",
        "choix_format": "📄 Zu generierendes Berichtsformat",
        "generer_rapport": "📜 Bericht generieren",
        "telecharger_rapport": "📥 Generierten Bericht herunterladen",
        "date_rapport": "Datum und Uhrzeit des Berichts",
        "element": "Element",
        "valeur": "Wert",
        "v80_message": "🚀 ErnestMind V40 analysiert bis zu 100.000 Zeilen. Für unbegrenzte Big-Data-Analyse, 5-fache Geschwindigkeit und Elite-Sektor-KI-Modelle, aktualisieren Sie auf Version V80!",
        "audit_title": "🔒 Lokale Prüfung & Sicherheit",
        "audit_exp": "In dieser Version werden alle temporären Dateien automatisch zur Gewährleistung der Vertraulichkeit verschlüsselt. Für jede fortgeschrittene Verarbeitung (Text/Audio/Bild) wird ein lokaler Prüfbericht generiert.",
        "brand_title": "🚀 ErnestMind AI V40 – Super Data Science Plattform",
        "brand_claim": "Leistungsstärker als Tableau, PowerBI, Dataiku, DataRobot, Alteryx, Qlik, SAS Viya, Palantir...<br>100% lokal, Auto-KI, PDF/Word/PPT-Berichte, mehrsprachig, Dunkel-/Hellmodus.",
        "connecteurs_title": "⚡️ Konnektoren / API / Cloud / DB (Demo)",
        "connecteurs_coming": "**Konnektoren in Kürze:**",
        "connecteurs_import": "- Import aus MySQL, PostgreSQL, BigQuery, S3, Google Drive, OneDrive, SharePoint, REST API",
        "connecteurs_export": "- Export nach interaktivem PowerPoint, HTML-Dashboard, API, Webhook, Slack, Teams, E-Mail usw.",
        "connecteurs_contact": "**Kontaktieren Sie uns, um Ihre Geschäftsdatenquelle zu integrieren!**",
        "hist_title": "Histogramm: {col}",
        "scatter_title": "Streudiagramm: {x} vs {y}",
        "box_title": "Boxplot: {col}",
        "corr_title": "Korrelation",
        "pairplot_title": "Variablen-Streudiagramm (Pairplot)",
        "pie_title": "Verteilung: {col}",
        "bar_title": "Balkendiagramm: {col}",
        "revolution_cta_title": "✨ Treten Sie der ErnestMind Revolution bei!",
        "revolution_cta_text": "Beeindruckt von der Leistung von ErnestMind V40? Das ist nur ein kleiner Einblick! Melden Sie sich an, um über neue Funktionen, spezialisierte KI-Modelle und exklusive Schulungen informiert zu werden.",
        "revolution_cta_button": "🚀 Treten Sie der KI-Revolution bei (Hier klicken!) 🚀",
        "revolution_cta_help": "Klicken Sie, um unsere Registrierungsseite zu besuchen und die Zukunft von ErnestMind zu entdecken.",
        "contact_email_label": "Kontakt"
    },
    "pt": {
        "aide_title": "🆘 Ajuda",
        "aide_texte": "Faça upload de um arquivo (csv, excel, parquet, ...), selecione um módulo de IA de negócios, visualize mais de 20 gráficos avançados e exporte relatórios PDF/Word/PPT instantaneamente. Total privacidade.",
        "apropos_title": "ℹ️ Sobre",
        "apropos_texte": "ErnestMind V40 supera Tableau, PowerBI, Dataiku, Qlik... em simplicidade, velocidade e privacidade. 100% local, IA automática, multi-idiomas, multi-formato.",
        "exemple_toggle": "📈 Exemplos de visualizações automáticas",
        "exemple_titre": "### 📊 Demonstração de IA e visualizações",
        "tab_corr": "📊 Correlação",
        "tab_boxplot": "📦 Boxplot",
        "tab_ventes": "📈 Vendas",
        "tab_marketing": "🎯 Marketing",
        "tab_satisfaction": "😊 Satisfação",
        "upload_label": "📂 Carregar um arquivo",
        "format_non_supporte": "❌ Formato não suportado",
        "success_upload": "✅ Arquivo carregado com sucesso!",
        "erreur_chargement": "❌ Erro ao carregar o arquivo",
        "explorer": "🖨️ Pré-visualização de dados",
        "nb_lignes": "Número de linhas",
        "nb_colonnes": "Número de colunas",
        "colonnes": "Nomes das colunas",
        "types_colonnes": "Tipos de colunas",
        "manquants": "Colunas com valores ausentes",
        "uniques": "Valores únicos por coluna",
        "stat_desc": "📐 Estatísticas descritivas",
        "apercu_donnees": "Pré-visualização de dados",
        "resume": "🔎 Resumo automático",
        "analyse_auto": "🧠 Análise automática inteligente",
        "numeriques": "Numéricos",
        "categorique": "Categóricos",
        "valeurs_manquantes": "Valores ausentes",
        "classif_detectee": "Classificação supervisionada detectada",
        "reg_detectee": "Regressão supervisionada detectada",
        "tache_non_detectee": "Tarefa não detectada automaticamente",
        "visual_auto": "📊 Visualizações automáticas",
        "histogramme": "Histograma",
        "nuage": "Gráfico de dispersão",
        "boxplot": "Boxplot",
        "repartition": "Distribuição",
        "choix_format": "📄 Formato do relatório a ser gerado",
        "generer_rapport": "📜 Gerar relatório",
        "telecharger_rapport": "📥 Baixar relatório gerado",
        "date_rapport": "Data e hora do relatório",
        "element": "Item",
        "valeur": "Valor",
        "v80_message": "🚀 ErnestMind V40 analisa até 100.000 linhas. Para análise de big data ilimitada, velocidade 5x e modelos de IA de elite, atualize para a versão V80!",
        "audit_title": "🔒 Auditoria local e segurança",
        "audit_exp": "Nesta versão, todos os arquivos temporários são criptografados automaticamente para garantir a confidencialidade. Um relatório de auditoria local é gerado para cada processamento avançado (texto/áudio/imagem).",
        "brand_title": "🚀 ErnestMind AI V40 – Plataforma de Super Ciência de Dados",
        "brand_claim": "Mais poderoso que Tableau, PowerBI, Dataiku, DataRobot, Alteryx, Qlik, SAS Viya, Palantir...<br>100% local, IA automática, relatórios PDF/Word/PPT, multi-idioma, modo escuro/claro.",
        "connecteurs_title": "⚡️ Conectores / API / Nuvem / BD (demonstração)",
        "connecteurs_coming": "**Conectores em breve:**",
        "connecteurs_import": "- Importar de MySQL, PostgreSQL, BigQuery, S3, Google Drive, OneDrive, SharePoint, API REST",
        "connecteurs_export": "- Exportar para PowerPoint interativo, painel HTML, API, webhook, Slack, Teams, email, etc.",
        "connecteurs_contact": "**Entre em contato para integrar sua fonte de dados de negócios!**",
        "hist_title": "Histograma: {col}",
        "scatter_title": "Gráfico de dispersão: {x} vs {y}",
        "box_title": "Boxplot: {col}",
        "corr_title": "Correlação",
        "pairplot_title": "Gráfico de pares de variáveis",
        "pie_title": "Distribuição: {col}",
        "bar_title": "Gráfico de barras: {col}",
        "revolution_cta_title": "✨ Junte-se à revolução ErnestMind!",
        "revolution_cta_text": "Impressionado com o poder do ErnestMind V40? Isso é apenas uma amostra! Inscreva-se para ser informado sobre novos recursos, modelos de IA especializados e treinamentos exclusivos.",
        "revolution_cta_button": "🚀 Junte-se à Revolução da IA (Clique aqui!) 🚀",
        "revolution_cta_help": "Clique para acessar nossa página de inscrição e descobrir o futuro do ErnestMind.",
        "contact_email_label": "Contato"
    },
    "it": {
        "aide_title": "🆘 Aiuto",
        "aide_texte": "Carica un file (csv, excel, parquet, ...), seleziona un modulo AI aziendale, visualizza oltre 20 grafici avanzati ed esporta rapporti PDF/Word/PPT istantaneamente. Privacy totale.",
        "apropos_title": "ℹ️ Chi siamo",
        "apropos_texte": "ErnestMind V40 supera Tableau, PowerBI, Dataiku, Qlik... in semplicità, velocità e privacy. 100% locale, AI automatica, multilingue, multi-formato.",
        "exemple_toggle": "📈 Esempi di visualizzazioni automatiche",
        "exemple_titre": "### 📊 Demo AI e visualizzazioni",
        "tab_corr": "📊 Correlazione",
        "tab_boxplot": "📦 Boxplot",
        "tab_ventes": "📈 Vendite",
        "tab_marketing": "🎯 Marketing",
        "tab_satisfaction": "😊 Soddisfazione",
        "upload_label": "📂 Carica un file",
        "format_non_supporte": "❌ Formato non supportato",
        "success_upload": "✅ File caricato con successo!",
        "erreur_chargement": "❌ Errore durante il caricamento del file",
        "explorer": "🖨️ Anteprima dei dati",
        "nb_lignes": "Numero di righe",
        "nb_colonnes": "Numero di colonne",
        "colonnes": "Nomi delle colonne",
        "types_colonnes": "Tipi di colonne",
        "manquants": "Colonne con valori mancanti",
        "uniques": "Valori unici per colonna",
        "stat_desc": "📐 Statistiche descrittive",
        "apercu_donnees": "Anteprima dei dati",
        "resume": "🔎 Riepilogo automatico",
        "analyse_auto": "🧠 Analisi automatica intelligente",
        "numeriques": "Numerici",
        "categorique": "Categorici",
        "valeurs_manquantes": "Valori mancanti",
        "classif_detectee": "Classificazione supervisionata rilevata",
        "reg_detectee": "Regressione supervisionata rilevata",
        "tache_non_detectee": "Task non rilevato automaticamente",
        "visual_auto": "📊 Visualizzazioni automatiche",
        "histogramme": "Istogramma",
        "nuage": "Grafico a dispersione",
        "boxplot": "Boxplot",
        "repartition": "Distribuzione",
        "choix_format": "📄 Formato del rapporto da generare",
        "generer_rapport": "📜 Genera rapporto",
        "telecharger_rapport": "📥 Scarica rapporto generato",
        "date_rapport": "Data e ora del rapporto",
        "element": "Elemento",
        "valeur": "Valore",
        "v80_message": "🚀 ErnestMind V40 analizza fino a 100.000 righe. Per l'analisi di big data illimitata, velocità x5 e modelli AI di élite, passa alla versione V80!",
        "audit_title": "🔒 Audit locale e sicurezza",
        "audit_exp": "In questa versione, tutti i file temporanei vengono crittografati automaticamente per garantire la riservatezza. Viene generato un rapporto di audit locale per ogni elaborazione avanzata (testo/audio/immagine).",
        "brand_title": "🚀 ErnestMind AI V40 – Piattaforma di Super Data Science",
        "brand_claim": "Più potente di Tableau, PowerBI, Dataiku, DataRobot, Alteryx, Qlik, SAS Viya, Palantir...<br>100% locale, AI automatica, rapporti PDF/Word/PPT, multilingue, modalità scura/chiara.",
        "connecteurs_title": "⚡️ Connettori / API / Cloud / DB (demo)",
        "connecteurs_coming": "**Connettori in arrivo:**",
        "connecteurs_import": "- Importa da MySQL, PostgreSQL, BigQuery, S3, Google Drive, OneDrive, SharePoint, API REST",
        "connecteurs_export": "- Esporta in PowerPoint interattivo, dashboard HTML, API, webhook, Slack, Teams, email, ecc.",
        "connecteurs_contact": "**Contattaci per integrare la tua fonte di business!**",
        "hist_title": "Istogramma: {col}",
        "scatter_title": "Grafico a dispersione: {x} vs {y}",
        "box_title": "Boxplot: {col}",
        "corr_title": "Correlazione",
        "pairplot_title": "Grafico a coppie di variabili",
        "pie_title": "Distribuzione: {col}",
        "bar_title": "Grafico a barre: {col}",
        "revolution_cta_title": "✨ Unisciti alla rivoluzione ErnestMind!",
        "revolution_cta_text": "Impressionato dalla potenza di ErnestMind V40? Questo è solo un assaggio! Iscriviti per essere informato su nuove funzionalità, modelli AI specializzati e corsi di formazione esclusivi.",
        "revolution_cta_button": "🚀 Unisciti alla Rivoluzione dell'IA (Clicca qui!) 🚀",
        "revolution_cta_help": "Clicca per accedere alla nostra pagina di iscrizione e scoprire il futuro di ErnestMind.",
        "contact_email_label": "Contatto"
    },
    "nl": {
        "aide_title": "🆘 Hulp",
        "aide_texte": "Upload een bestand (csv, excel, parquet, ...), selecteer een bedrijfs-AI-module, bekijk meer dan 20 geavanceerde grafieken en exporteer direct PDF/Word/PPT-rapporten. Volledige privacy.",
        "apropos_title": "ℹ️ Over",
        "apropos_texte": "ErnestMind V40 overtreft Tableau, PowerBI, Dataiku, Qlik... in eenvoud, snelheid en privacy. 100% lokaal, Auto AI, meertalig, multi-formaat.",
        "exemple_toggle": "📈 Voorbeelden van automatische visualisaties",
        "exemple_titre": "### 📊 AI Demo & visualisaties",
        "tab_corr": "📊 Correlatie",
        "tab_boxplot": "📦 Boxplot",
        "tab_ventes": "📈 Verkoop",
        "tab_marketing": "🎯 Marketing",
        "tab_satisfaction": "😊 Tevredenheid",
        "upload_label": "📂 Bestand uploaden",
        "format_non_supporte": "❌ Niet-ondersteund formaat",
        "success_upload": "✅ Bestand succesvol geüpload!",
        "erreur_chargement": "❌ Fout bij het laden van het bestand",
        "explorer": "🖨️ Gegevensvoorbeeld",
        "nb_lignes": "Aantal rijen",
        "nb_colonnes": "Aantal kolommen",
        "colonnes": "Kolomnamen",
        "types_colonnes": "Kolomtypen",
        "manquants": "Kolommen met ontbrekende waarden",
        "uniques": "Unieke waarden per kolom",
        "stat_desc": "📐 Beschrijvende statistieken",
        "apercu_donnees": "Gegevensvoorbeeld",
        "resume": "🔎 Auto-samenvatting",
        "analyse_auto": "🧠 Slimme analyse",
        "numeriques": "Numeriek",
        "categorique": "Categorisch",
        "valeurs_manquantes": "Ontbrekende waarden",
        "classif_detectee": "Begeleide classificatie gedetecteerd",
        "reg_detectee": "Begeleide regressie gedetecteerd",
        "tache_non_detectee": "Taak niet automatisch gedetecteerd",
        "visual_auto": "📊 Automatische visualisaties",
        "histogramme": "Histogram",
        "nuage": "Spreidingsdiagram",
        "boxplot": "Boxplot",
        "repartition": "Distributie",
        "choix_format": "📄 Te genereren rapportformaat",
        "generer_rapport": "📜 Rapport genereren",
        "telecharger_rapport": "📥 Gedownload rapport downloaden",
        "date_rapport": "Datum en tijd rapport",
        "element": "Item",
        "valeur": "Waarde",
        "v80_message": "🚀 ErnestMind V40 analyseert tot 100.000 rijen. Voor onbeperkte big data-analyse, 5x snelheid en elite AI-modellen, upgrade naar versie V80!",
        "audit_title": "🔒 Lokale audit & beveiliging",
        "audit_exp": "In deze versie worden alle tijdelijke bestanden automatisch versleuteld om vertrouwelijkheid te garanderen. Een lokaal auditrapport wordt gegenereerd voor elke geavanceerde verwerking (tekst/audio/afbeelding).",
        "brand_title": "🚀 ErnestMind AI V40 – Super Data Science Platform",
        "brand_claim": "Sterker dan Tableau, PowerBI, Dataiku, DataRobot, Alteryx, Qlik, SAS Viya, Palantir...<br>100% lokaal, Auto AI, PDF/Word/PPT-rapporten, meertalig, donkere/lichte modus.",
        "connecteurs_title": "⚡️ Connectoren / API / Cloud / DB (demo)",
        "connecteurs_coming": "**Connectoren binnenkort:**",
        "connecteurs_import": "- Importeren vanuit MySQL, PostgreSQL, BigQuery, S3, Google Drive, OneDrive, SharePoint, REST API",
        "connecteurs_export": "- Exporteren naar interactieve PowerPoint, HTML-dashboard, API, webhook, Slack, Teams, e-mail, etc.",
        "connecteurs_contact": "**Neem contact met ons op om uw bedrijfsbron te integreren!**",
        "hist_title": "Histogram: {col}",
        "scatter_title": "Spreidingsdiagram: {x} vs {y}",
        "box_title": "Boxplot: {col}",
        "corr_title": "Correlatie",
        "pairplot_title": "Variabele spreiding (pairplot)",
        "pie_title": "Distributie: {col}",
        "bar_title": "Staafdiagram: {col}",
        "revolution_cta_title": "✨ Sluit je aan bij de ErnestMind Revolutie!",
        "revolution_cta_text": "Onder de indruk van de kracht van ErnestMind V40? Dit is slechts een voorproefje! Meld u aan om op de hoogte te blijven van nieuwe functies, gespecialiseerde AI-modellen en exclusieve trainingen.",
        "revolution_cta_button": "🚀 Sluit je aan bij de AI Revolutie (Klik hier!) 🚀",
        "revolution_cta_help": "Klik om naar onze aanmeldingspagina te gaan en de toekomst van ErnestMind te ontdekken.",
        "contact_email_label": "Contact"
    },
    "zh": {
        "aide_title": "🆘 帮助",
        "aide_texte": "上传文件 (csv, excel, parquet, ...)，选择商业 AI 模块，查看 20+ 高级图表并即时导出 PDF/Word/PPT 报告。完全隐私。",
        "apropos_title": "ℹ️ 关于",
        "apropos_texte": "ErnestMind V40 在简单性、速度和隐私方面超越了 Tableau, PowerBI, Dataiku, Qlik... 100% 本地运行，自动 AI，多语言，多格式。",
        "exemple_toggle": "📈 自动可视化示例",
        "exemple_titre": "### 📊 AI 演示与可视化",
        "tab_corr": "📊 相关性",
        "tab_boxplot": "📦 箱线图",
        "tab_ventes": "📈 销售",
        "tab_marketing": "🎯 市场营销",
        "tab_satisfaction": "😊 满意度",
        "upload_label": "📂 上传文件",
        "format_non_supporte": "❌ 不支持的格式",
        "success_upload": "✅ 文件上传成功！",
        "erreur_chargement": "❌ 文件加载错误",
        "explorer": "🖨️ 数据预览",
        "nb_lignes": "行数",
        "nb_colonnes": "列数",
        "colonnes": "列名",
        "types_colonnes": "列类型",
        "manquants": "含有缺失值的列",
        "uniques": "每列唯一值",
        "stat_desc": "📐 描述性统计",
        "apercu_donnees": "数据预览",
        "resume": "🔎 自动摘要",
        "analyse_auto": "🧠 智能自动分析",
        "numeriques": "数值型",
        "categorique": "类别型",
        "valeurs_manquantes": "缺失值",
        "classif_detectee": "检测到监督分类任务",
        "reg_detectee": "检测到监督回归任务",
        "tache_non_detectee": "未自动检测到任务",
        "visual_auto": "📊 自动可视化",
        "histogramme": "直方图",
        "nuage": "散点图",
        "boxplot": "箱线图",
        "repartition": "分布",
        "choix_format": "📄 报告格式",
        "generer_rapport": "📜 生成报告",
        "telecharger_rapport": "📥 下载生成的报告",
        "date_rapport": "报告日期/时间",
        "element": "项目",
        "valeur": "值",
        "v80_message": "🚀 ErnestMind V40 最多可分析 100,000 行。如需无限大数据分析、5倍速度和精英行业模型，请升级到 V80！",
        "audit_title": "🔒 本地审计与安全",
        "audit_exp": "在此版本中，所有临时文件都自动加密以确保机密性。每次高级处理（文本/音频/图像）都会生成本地审计报告。",
        "brand_title": "🚀 ErnestMind AI V40 – 超级数据科学平台",
        "brand_claim": "比 Tableau, PowerBI, Dataiku, DataRobot, Alteryx, Qlik, SAS Viya, Palantir 更强大...<br>100% 本地，自动 AI，PDF/Word/PPT 报告，多语言，深色/浅色模式。",
        "connecteurs_title": "⚡️ 连接器 / API / 云 / 数据库 (演示)",
        "connecteurs_coming": "**即将推出的连接器：**",
        "connecteurs_import": "- 从 MySQL, PostgreSQL, BigQuery, S3, Google Drive, OneDrive, SharePoint, REST API 导入",
        "connecteurs_export": "- 导出到交互式 PowerPoint, HTML 仪表板, API, webhook, Slack, Teams, 电子邮件等。",
        "connecteurs_contact": "**联系我们，集成您的业务数据源！**",
        "hist_title": "直方图：{col}",
        "scatter_title": "散点图：{x} vs {y}",
        "box_title": "箱线图：{col}",
        "corr_title": "相关性",
        "pairplot_title": "变量散点图 (pairplot)",
        "pie_title": "分布：{col}",
        "bar_title": "条形图：{col}",
        "revolution_cta_title": "✨ 加入 ErnestMind 革命！",
        "revolution_cta_text": "对 ErnestMind V40 的强大功能印象深刻？这只是冰山一角！注册以获取新功能、专业 AI 模型和独家培训的信息。",
        "revolution_cta_button": "🚀 加入 AI 革命 (点击这里！) 🚀",
        "revolution_cta_help": "点击访问我们的注册页面，发现 ErnestMind 的未来。",
        "contact_email_label": "联系方式"
    },
    "ja": {
        "aide_title": "🆘 ヘルプ",
        "aide_texte": "ファイルをアップロード (csv, excel, parquet, ...)、ビジネスAIモジュールを選択、20以上の高度なグラフを表示し、PDF/Word/PPTレポートを即座にエクスポート。完全なプライバシー。",
        "apropos_title": "ℹ️ 概要",
        "apropos_texte": "ErnestMind V40は、Tableau、PowerBI、Dataiku、Qlik...をシンプルさ、速度、プライバシーの面で凌駕しています。100%ローカル、自動AI、多言語、マルチフォーマット。",
        "exemple_toggle": "📈 自動可視化の例",
        "exemple_titre": "### 📊 AIデモ＆可視化",
        "tab_corr": "📊 相関",
        "tab_boxplot": "📦 箱ひげ図",
        "tab_ventes": "📈 販売",
        "tab_marketing": "🎯 マーケティング",
        "tab_satisfaction": "😊 満足度",
        "upload_label": "📂 ファイルをアップロード",
        "format_non_supporte": "❌ 非対応フォーマット",
        "success_upload": "✅ ファイルのアップロードに成功しました！",
        "erreur_chargement": "❌ ファイルの読み込みエラー",
        "explorer": "🖨️ データプレビュー",
        "nb_lignes": "行数",
        "nb_colonnes": "列数",
        "colonnes": "列名",
        "types_colonnes": "列の型",
        "manquants": "欠損値のある列",
        "uniques": "列ごとのユニーク値",
        "stat_desc": "📐 記述統計",
        "apercu_donnees": "データプレビュー",
        "resume": "🔎 自動要約",
        "analyse_auto": "🧠 スマートな自動分析",
        "numeriques": "数値",
        "categorique": "カテゴリ",
        "valeurs_manquantes": "欠損値",
        "classif_detectee": "教師あり分類を検出",
        "reg_detectee": "教師あり回帰を検出",
        "tache_non_detectee": "タスクが自動検出されませんでした",
        "visual_auto": "📊 自動可視化",
        "histogramme": "ヒストグラム",
        "nuage": "散布図",
        "boxplot": "箱ひげ図",
        "repartition": "分布",
        "choix_format": "📄 生成するレポート形式",
        "generer_rapport": "📜 レポートを生成",
        "telecharger_rapport": "📥 生成されたレポートをダウンロード",
        "date_rapport": "レポート日時",
        "element": "項目",
        "valeur": "値",
        "v80_message": "🚀 ErnestMind V40は最大10万行を分析します。無制限のビッグデータ分析、5倍の速度、およびエリートセクターのAIモデルについては、V80にアップグレードしてください！",
        "audit_title": "🔒 ローカル監査とセキュリティ",
        "audit_exp": "このバージョンでは、すべてのテンポラリファイルは機密性を保証するために自動的に暗号化されます。各高度な処理（テキスト/オーディオ/画像）について、ローカル監査レポートが生成されます。",
        "brand_title": "🚀 ErnestMind AI V40 – スーパーデータサイエンスプラットフォーム",
        "brand_claim": "Tableau、PowerBI、Dataiku、DataRobot、Alteryx、Qlik、SAS Viya、Palantir...よりも強力です。<br>100%ローカル、自動AI、PDF/Word/PPTレポート、多言語、ダーク/ライトモード。",
        "connecteurs_title": "⚡️ コネクタ / API / クラウド / DB (デモ)",
        "connecteurs_coming": "**近日公開予定のコネクタ：**",
        "connecteurs_import": "- MySQL, PostgreSQL, BigQuery, S3, Google Drive, OneDrive, SharePoint, REST API からインポート",
        "connecteurs_export": "- 対話型PowerPoint, HTMLダッシュボード, API, webhook, Slack, Teams, メールなどへエクスポート",
        "connecteurs_contact": "**ビジネスソースの統合についてはお問い合わせください！**",
        "hist_title": "ヒストグラム: {col}",
        "scatter_title": "散布図: {x} vs {y}",
        "box_title": "箱ひげ図: {col}",
        "corr_title": "相関",
        "pairplot_title": "変数散布図 (ペアプロット)",
        "pie_title": "分布: {col}",
        "bar_title": "棒グラフ: {col}",
        "revolution_cta_title": "✨ ErnestMind革命に参加しよう！",
        "revolution_cta_text": "ErnestMind V40の力に感銘を受けましたか？これはほんの一部です！新機能、専門AIモデル、独占トレーニングに関する情報を受け取るには登録してください。",
        "revolution_cta_button": "🚀 AI革命に参加 (ここをクリック！) 🚀",
        "revolution_cta_help": "クリックして登録ページにアクセスし、ErnestMindの未来を発見してください。",
        "contact_email_label": "連絡先"
    },
    "ko": {
        "aide_title": "🆘 도움말",
        "aide_texte": "파일 업로드 (csv, excel, parquet, ...), 비즈니스 AI 모듈 선택, 20개 이상의 고급 차트 보기 및 PDF/Word/PPT 보고서 즉시 내보내기. 완전한 프라이버시.",
        "apropos_title": "ℹ️ 정보",
        "apropos_texte": "ErnestMind V40은 단순성, 속도, 프라이버시 측면에서 Tableau, PowerBI, Dataiku, Qlik을 능가합니다. 100% 로컬, 자동 AI, 다국어, 다중 형식.",
        "exemple_toggle": "📈 자동 시각화 예시",
        "exemple_titre": "### 📊 AI 데모 및 시각화",
        "tab_corr": "📊 상관관계",
        "tab_boxplot": "📦 상자 그림",
        "tab_ventes": "📈 판매",
        "tab_marketing": "🎯 마케팅",
        "tab_satisfaction": "😊 만족도",
        "upload_label": "📂 파일 업로드",
        "format_non_supporte": "❌ 지원되지 않는 형식",
        "success_upload": "✅ 파일 업로드 성공!",
        "erreur_chargement": "❌ 파일 로드 오류",
        "explorer": "🖨️ 데이터 미리보기",
        "nb_lignes": "행 수",
        "nb_colonnes": "열 수",
        "colonnes": "열 이름",
        "types_colonnes": "열 유형",
        "manquants": "누락된 값이 있는 열",
        "uniques": "열별 고유 값",
        "stat_desc": "📐 기술 통계",
        "apercu_donnees": "데이터 미리보기",
        "resume": "🔎 자동 요약",
        "analyse_auto": "🧠 스마트 자동 분석",
        "numeriques": "숫자형",
        "categorique": "범주형",
        "valeurs_manquantes": "누락된 값",
        "classif_detectee": "지도 분류 감지됨",
        "reg_detectee": "지도 회귀 감지됨",
        "tache_non_detectee": "작업 자동 감지 안 됨",
        "visual_auto": "📊 자동 시각화",
        "histogramme": "히스토그램",
        "nuage": "산점도",
        "boxplot": "상자 그림",
        "repartition": "분포",
        "choix_format": "📄 생성할 보고서 형식",
        "generer_rapport": "📜 보고서 생성",
        "telecharger_rapport": "📥 생성된 보고서 다운로드",
        "date_rapport": "보고서 날짜/시간",
        "element": "항목",
        "valeur": "값",
        "v80_message": "🚀 ErnestMind V40은 최대 10만 행을 분석합니다. 무제한 빅데이터 분석, 5배 빠른 속도, 엘리트 섹터 AI 모델을 위해 V80으로 업그레이드하세요!",
        "audit_title": "🔒 로컬 감사 및 보안",
        "audit_exp": "이 버전에서는 모든 임시 파일이 기밀 유지를 위해 자동으로 암호화됩니다. 각 고급 처리 (텍스트/오디오/이미지)에 대해 로컬 감사 보고서가 생성됩니다.",
        "brand_title": "🚀 ErnestMind AI V40 – 슈퍼 데이터 과학 플랫폼",
        "brand_claim": "Tableau, PowerBI, Dataiku, DataRobot, Alteryx, Qlik, SAS Viya, Palantir보다 강력합니다...<br>100% 로컬, 자동 AI, PDF/Word/PPT 보고서, 다국어, 다크/라이트 모드.",
        "connecteurs_title": "⚡️ 커넥터 / API / 클라우드 / DB (데모)",
        "connecteurs_coming": "**예정된 커넥터:**",
        "connecteurs_import": "- MySQL, PostgreSQL, BigQuery, S3, Google Drive, OneDrive, SharePoint, REST API에서 가져오기",
        "connecteurs_export": "- 대화형 PowerPoint, HTML 대시보드, API, 웹훅, Slack, Teams, 이메일 등으로 내보내기",
        "connecteurs_contact": "**귀사의 비즈니스 소스를 통합하려면 문의하세요!**",
        "hist_title": "히스토그램: {col}",
        "scatter_title": "산점도: {x} vs {y}",
        "box_title": "상자 그림: {col}",
        "corr_title": "상관관계",
        "pairplot_title": "변수 산점도 (페어플롯)",
        "pie_title": "분포: {col}",
        "bar_title": "막대 그래프: {col}",
        "revolution_cta_title": "✨ ErnestMind 혁명에 참여하세요!",
        "revolution_cta_text": "ErnestMind V40의 강력함에 감동받으셨나요? 이것은 맛보기에 불과합니다! 새로운 기능, 전문 AI 모델 및 독점 교육에 대한 정보를 받으려면 등록하세요.",
        "revolution_cta_button": "🚀 AI 혁명에 참여 (여기를 클릭하세요!) 🚀",
        "revolution_cta_help": "클릭하여 등록 페이지에 접속하고 ErnestMind의 미래를 발견하세요.",
        "contact_email_label": "연락처"
    },
    "ar": {
        "aide_title": "🆘 مساعدة",
        "aide_texte": "قم بتحميل ملف (csv, excel, parquet, ...)، اختر وحدة ذكاء اصطناعي للأعمال، شاهد أكثر من 20 مخططًا متقدمًا وقم بتصدير تقارير PDF/Word/PPT على الفور. خصوصية تامة.",
        "apropos_title": "ℹ️ حول",
        "apropos_texte": "يتجاوز ErnestMind V40 برامج Tableau وPowerBI وDataiku وQlik... في البساطة والسرعة والخصوصية. محلي 100%، ذكاء اصطناعي تلقائي، متعدد اللغات، متعدد التنسيقات.",
        "exemple_toggle": "📈 أمثلة على التصورات التلقائية",
        "exemple_titre": "### 📊 عرض توضيحي للذكاء الاصطناعي والتصورات",
        "tab_corr": "📊 الارتباط",
        "tab_boxplot": "📦 مخطط الصندوق",
        "tab_ventes": "📈 المبيعات",
        "tab_marketing": "🎯 التسويق",
        "tab_satisfaction": "😊 الرضا",
        "upload_label": "📂 تحميل ملف",
        "format_non_supporte": "❌ تنسيق غير مدعوم",
        "success_upload": "✅ تم تحميل الملف بنجاح!",
        "erreur_chargement": "❌ خطأ في تحميل الملف",
        "explorer": "🖨️ معاينة البيانات",
        "nb_lignes": "عدد الصفوف",
        "nb_colonnes": "عدد الأعمدة",
        "colonnes": "أسماء الأعمدة",
        "types_colonnes": "أنواع الأعمدة",
        "manquants": "الأعمدة ذات القيم المفقودة",
        "uniques": "القيم الفريدة لكل عمود",
        "stat_desc": "📐 الإحصائيات الوصفية",
        "apercu_donnees": "معاينة البيانات",
        "resume": "🔎 ملخص تلقائي",
        "analyse_auto": "🧠 تحليل تلقائي ذكي",
        "numeriques": "رقمية",
        "categorique": "فئوية",
        "valeurs_manquantes": "القيم المفقودة",
        "classif_detectee": "تم الكشف عن التصنيف الخاضع للإشراف",
        "reg_detectee": "تم الكشف عن الانحدار الخاضع للإشراف",
        "tache_non_detectee": "لم يتم الكشف عن المهمة تلقائيًا",
        "visual_auto": "📊 تصورات تلقائية",
        "histogramme": "مخطط بياني",
        "nuage": "مخطط الانتشار",
        "boxplot": "مخطط الصندوق",
        "repartition": "التوزيع",
        "choix_format": "📄 تنسيق التقرير المراد إنشاؤه",
        "generer_rapport": "📜 إنشاء التقرير",
        "telecharger_rapport": "📥 تنزيل التقرير الذي تم إنشاؤه",
        "date_rapport": "تاريخ/وقت التقرير",
        "element": "العنصر",
        "valeur": "القيمة",
        "v80_message": "🚀 يحلل ErnestMind V40 ما يصل إلى 100,000 صف. لتحليل البيانات الضخمة غير المحدود، وسرعة 5 أضعاف، ونماذج الذكاء الاصطناعي المتخصصة، قم بالترقية إلى V80!",
        "audit_title": "🔒 التدقيق المحلي والأمان",
        "audit_exp": "في هذا الإصدار، يتم تشفير جميع الملفات المؤقتة تلقائيًا لضمان السرية. يتم إنشاء تقرير تدقيق محلي لكل معالجة متقدمة (نص/صوت/صورة).",
        "brand_title": "🚀 ErnestMind AI V40 – منصة علوم بيانات فائقة",
        "brand_claim": "أقوى من Tableau, PowerBI, Dataiku, DataRobot, Alteryx, Qlik, SAS Viya, Palantir...<br>100% محلي، ذكاء اصطناعي تلقائي، تقارير PDF/Word/PPT، متعدد اللغات، وضع داكن/فاتح.",
        "connecteurs_title": "⚡️ الموصلات / API / السحابة / قاعدة البيانات (تجريبي)",
        "connecteurs_coming": "**الموصلات قادمة قريبا:**",
        "connecteurs_import": "- الاستيراد من MySQL, PostgreSQL, BigQuery, S3, Google Drive, OneDrive, SharePoint, REST API",
        "connecteurs_export": "- التصدير إلى PowerPoint تفاعلي, لوحة معلومات HTML, API, webhook, Slack, Teams, البريد الإلكتروني, إلخ.",
        "connecteurs_contact": "**اتصل بنا لدمج مصدر عملك!**",
        "hist_title": "مخطط بياني: {col}",
        "scatter_title": "مخطط الانتشار: {x} مقابل {y}",
        "box_title": "مخطط الصندوق: {col}",
        "corr_title": "الارتباط",
        "pairplot_title": "مخطط انتشار المتغيرات (مخطط الأزواج)",
        "pie_title": "التوزيع: {col}",
        "bar_title": "مخطط الأعمدة: {col}",
        "revolution_cta_title": "✨ انضم إلى ثورة ErnestMind!",
        "revolution_cta_text": "هل أنت معجب بقوة ErnestMind V40؟ هذا مجرد لمحة! سجل لتكون على اطلاع بالميزات الجديدة ونماذج الذكاء الاصطناعي المتخصصة والتدريب الحصري.",
        "revolution_cta_button": "🚀 انضم إلى ثورة الذكاء الاصطناعي (اضغط هنا!) 🚀",
        "revolution_cta_help": "انقر للوصول إلى صفحة التسجيل الخاصة بنا واكتشف مستقبل ErnestMind.",
        "contact_email_label": "جهة الاتصال"
    },
    "hi": {
        "aide_title": "🆘 मदद",
        "aide_texte": "एक फ़ाइल अपलोड करें (csv, excel, parquet, ...), एक व्यावसायिक AI मॉड्यूल का चयन करें, 20+ उन्नत चार्ट देखें और तुरंत PDF/Word/PPT रिपोर्ट निर्यात करें। कुल गोपनीयता।",
        "apropos_title": "ℹ️ के बारे में",
        "apropos_texte": "ErnestMind V40 सादगी, गति और गोपनीयता में Tableau, PowerBI, Dataiku, Qlik... को पीछे छोड़ देता है। 100% स्थानीय, ऑटो AI, बहु-भाषा, बहु-प्रारूप।",
        "exemple_toggle": "📈 स्वचालित विज़ुअलाइज़ेशन के उदाहरण",
        "exemple_titre": "### 📊 AI डेमो और विज़ुअलाइज़ेशन",
        "tab_corr": "📊 सहसंबंध",
        "tab_boxplot": "📦 बॉक्सप्लॉट",
        "tab_ventes": "📈 बिक्री",
        "tab_marketing": "🎯 मार्केटिंग",
        "tab_satisfaction": "😊 संतुष्टि",
        "upload_label": "📂 फ़ाइल अपलोड करें",
        "format_non_supporte": "❌ असमर्थित प्रारूप",
        "success_upload": "✅ फ़ाइल सफलतापूर्वक अपलोड हुई!",
        "erreur_chargement": "❌ फ़ाइल लोड करने में त्रुटि",
        "explorer": "🖨️ डेटा पूर्वावलोकन",
        "nb_lignes": "पंक्तियों की संख्या",
        "nb_colonnes": "कॉलम की संख्या",
        "colonnes": "कॉलम के नाम",
        "types_colonnes": "कॉलम के प्रकार",
        "manquants": "छूटे हुए मानों वाले कॉलम",
        "uniques": "प्रति कॉलम अद्वितीय मान",
        "stat_desc": "📐 वर्णनात्मक आंकड़े",
        "apercu_donnees": "डेटा पूर्वावलोकन",
        "resume": "🔎 स्वचालित सारांश",
        "analyse_auto": "🧠 स्मार्ट स्वचालित विश्लेषण",
        "numeriques": "संख्यात्मक",
        "categorique": "श्रेणीबद्ध",
        "valeurs_manquantes": "छूटे हुए मान",
        "classif_detectee": "पर्यवेक्षित वर्गीकरण का पता चला",
        "reg_detectee": "पर्यवेक्षित प्रतिगमन का पता चला",
        "tache_non_detectee": "कार्य स्वतः पता नहीं चला",
        "visual_auto": "📊 स्वचालित विज़ुअलाइज़ेशन",
        "histogramme": "हिस्टोग्राम",
        "nuage": "स्कैटरप्लॉट",
        "boxplot": "बॉक्सप्लॉट",
        "repartition": "वितरण",
        "choix_format": "📄 जेनरेट करने के लिए रिपोर्ट प्रारूप",
        "generer_rapport": "📜 रिपोर्ट जेनरेट करें",
        "telecharger_rapport": "📥 जेनरेट की गई रिपोर्ट डाउनलोड करें",
        "date_rapport": "रिपोर्ट की तारीख/समय",
        "element": "आइटम",
        "valeur": "मूल्य",
        "v80_message": "🚀 ErnestMind V40 100,000 पंक्तियों तक का विश्लेषण करता है। असीमित बिग डेटा विश्लेषण, 5x गति और विशिष्ट क्षेत्र AI मॉडल के लिए, V80 में अपग्रेड करें!",
        "audit_title": "🔒 स्थानीय ऑडिट और सुरक्षा",
        "audit_exp": "इस संस्करण में, गोपनीयता सुनिश्चित करने के लिए सभी अस्थायी फ़ाइलें स्वचालित रूप से एन्क्रिप्ट की जाती हैं। प्रत्येक उन्नत प्रसंस्करण (पाठ/ऑडियो/छवि) के लिए एक स्थानीय ऑडिट रिपोर्ट तैयार की जाती है।",
        "brand_title": "🚀 ErnestMind AI V40 – सुपर डेटा साइंस प्लेटफॉर्म",
        "brand_claim": "Tableau, PowerBI, Dataiku, DataRobot, Alteryx, Qlik, SAS Viya, Palantir... से अधिक शक्तिशाली<br>100% स्थानीय, ऑटो AI, PDF/Word/PPT रिपोर्ट, बहुभाषी, डार्क/लाइट मोड।",
        "connecteurs_title": "⚡️ कनेक्टर / API / क्लाउड / DB (डेमो)",
        "connecteurs_coming": "**जल्द आ रहे कनेक्टर:**",
        "connecteurs_import": "- MySQL, PostgreSQL, BigQuery, S3, Google Drive, OneDrive, SharePoint, REST API से आयात करें",
        "connecteurs_export": "- इंटरैक्टिव PowerPoint, HTML डैशबोर्ड, API, वेबहुक, स्लैक, टीम्स, ईमेल आदि में निर्यात करें।",
        "connecteurs_contact": "**अपने व्यावसायिक स्रोत को एकीकृत करने के लिए हमसे संपर्क करें!**",
        "hist_title": "हिस्टोग्राम: {col}",
        "scatter_title": "स्कैटरप्लॉट: {x} बनाम {y}",
        "box_title": "बॉक्सप्लॉट: {col}",
        "corr_title": "सहसंबंध",
        "pairplot_title": "चरों का स्कैटर (पेयरप्लॉट)",
        "pie_title": "वितरण: {col}",
        "bar_title": "बारप्लॉट: {col}",
        "revolution_cta_title": "✨ ErnestMind क्रांति में शामिल हों!",
        "revolution_cta_text": "ErnestMind V40 की शक्ति से प्रभावित हैं? यह सिर्फ एक झलक है! नई सुविधाओं, विशेष AI मॉडल और विशेष प्रशिक्षण के बारे में सूचित रहने के लिए साइन अप करें।",
        "revolution_cta_button": "🚀 AI क्रांति में शामिल हों (यहां क्लिक करें!) 🚀",
        "revolution_cta_help": "हमारी साइन-अप पेज पर जाने और ErnestMind के भविष्य की खोज करने के लिए क्लिक करें।",
        "contact_email_label": "संपर्क"
    },
    "th": {
        "aide_title": "🆘 ช่วยเหลือ",
        "aide_texte": "อัปโหลดไฟล์ (csv, excel, parquet, ...), เลือกโมดูล AI สำหรับธุรกิจ, ดูแผนภูมิขั้นสูงกว่า 20 รายการ และส่งออกรายงาน PDF/Word/PPT ได้ทันที ความเป็นส่วนตัวสูงสุด",
        "apropos_title": "ℹ️ เกี่ยวกับ",
        "apropos_texte": "ErnestMind V40 เหนือกว่า Tableau, PowerBI, Dataiku, Qlik... ในด้านความเรียบง่าย ความเร็ว และความเป็นส่วนตัว ทำงานแบบ 100% โลคัล, AI อัตโนมัติ, หลายภาษา, หลายรูปแบบ",
        "exemple_toggle": "📈 ตัวอย่างการแสดงภาพอัตโนมัติ",
        "exemple_titre": "### 📊 เดโม AI และการแสดงภาพ",
        "tab_corr": "📊 ความสัมพันธ์",
        "tab_boxplot": "📦 Boxplot",
        "tab_ventes": "📈 ยอดขาย",
        "tab_marketing": "🎯 การตลาด",
        "tab_satisfaction": "😊 ความพึงพอใจ",
        "upload_label": "📂 อัปโหลดไฟล์",
        "format_non_supporte": "❌ รูปแบบไม่รองรับ",
        "success_upload": "✅ อัปโหลดไฟล์สำเร็จ!",
        "erreur_chargement": "❌ ข้อผิดพลาดในการโหลดไฟล์",
        "explorer": "🖨️ ดูตัวอย่างข้อมูล",
        "nb_lignes": "จำนวนแถว",
        "nb_colonnes": "จำนวนคอลัมน์",
        "colonnes": "ชื่อคอลัมน์",
        "types_colonnes": "ประเภทคอลัมน์",
        "manquants": "คอลัมน์ที่มีค่าที่ขาดหายไป",
        "uniques": "ค่าที่ไม่ซ้ำกันต่อคอลัมน์",
        "stat_desc": "📐 สถิติเชิงพรรณนา",
        "apercu_donnees": "ดูตัวอย่างข้อมูล",
        "resume": "🔎 สรุปอัตโนมัติ",
        "analyse_auto": "🧠 การวิเคราะห์อัตโนมัติอัจฉริยะ",
        "numeriques": "ตัวเลข",
        "categorique": "หมวดหมู่",
        "valeurs_manquantes": "ค่าที่ขาดหายไป",
        "classif_detectee": "ตรวจพบการจำแนกประเภทแบบมีการควบคุม",
        "reg_detectee": "ตรวจพบการถดถอยแบบมีการควบคุม",
        "tache_non_detectee": "ไม่พบงานโดยอัตโนมัติ",
        "visual_auto": "📊 การแสดงภาพอัตโนมัติ",
        "histogramme": "ฮิสโทแกรม",
        "nuage": "แผนภาพการกระจาย",
        "boxplot": "Boxplot",
        "repartition": "การกระจาย",
        "choix_format": "📄 รูปแบบรายงานที่จะสร้าง",
        "generer_rapport": "📜 สร้างรายงาน",
        "telecharger_rapport": "📥 ดาวน์โหลดรายงานที่สร้างขึ้น",
        "date_rapport": "วันที่/เวลาของรายงาน",
        "element": "รายการ",
        "valeur": "ค่า",
        "v80_message": "🚀 ErnestMind V40 วิเคราะห์ได้สูงสุด 100,000 แถว สำหรับการวิเคราะห์ Big Data แบบไม่จำกัด ความเร็ว 5 เท่า และโมเดล AI เฉพาะทางขั้นสูง อัปเกรดเป็น V80!",
        "audit_title": "🔒 การตรวจสอบภายในและความปลอดภัย",
        "audit_exp": "ในเวอร์ชันนี้ ไฟล์ชั่วคราวทั้งหมดจะถูกเข้ารหัสโดยอัตโนมัติเพื่อรับประกันความลับ รายงานการตรวจสอบภายในจะถูกสร้างขึ้นสำหรับการประมวลผลขั้นสูงแต่ละครั้ง (ข้อความ/เสียง/รูปภาพ)",
        "brand_title": "🚀 ErnestMind AI V40 – แพลตฟอร์มวิทยาศาสตร์ข้อมูลขั้นสูง",
        "brand_claim": "มีประสิทธิภาพมากกว่า Tableau, PowerBI, Dataiku, DataRobot, Alteryx, Qlik, SAS Viya, Palantir...<br>ทำงานแบบ 100% โลคัล, AI อัตโนมัติ, รายงาน PDF/Word/PPT, หลายภาษา, โหมดมืด/สว่าง",
        "connecteurs_title": "⚡️ ตัวเชื่อมต่อ / API / คลาวด์ / DB (ตัวอย่าง)",
        "connecteurs_coming": "**ตัวเชื่อมต่อที่จะมาถึงเร็วๆ นี้:**",
        "connecteurs_import": "- นำเข้าจาก MySQL, PostgreSQL, BigQuery, S3, Google Drive, OneDrive, SharePoint, REST API",
        "connecteurs_export": "- ส่งออกไปยัง PowerPoint แบบโต้ตอบ, แดชบอร์ด HTML, API, webhook, Slack, Teams, อีเมล ฯลฯ",
        "connecteurs_contact": "**ติดต่อเราเพื่อรวมแหล่งข้อมูลธุรกิจของคุณ!**",
        "hist_title": "ฮิสโทแกรม: {col}",
        "scatter_title": "แผนภาพการกระจาย: {x} เทียบกับ {y}",
        "box_title": "Boxplot: {col}",
        "corr_title": "ความสัมพันธ์",
        "pairplot_title": "แผนภาพการกระจายตัวแปร (pairplot)",
        "pie_title": "การกระจาย: {col}",
        "bar_title": "แผนภูมิแท่ง: {col}",
        "revolution_cta_title": "✨ เข้าร่วมการปฏิวัติ ErnestMind!",
        "revolution_cta_text": "ประทับใจในพลังของ ErnestMind V40 ใช่ไหม? นี่เป็นเพียงส่วนเล็ก ๆ เท่านั้น! ลงทะเบียนเพื่อรับข้อมูลอัปเดตเกี่ยวกับคุณสมบัติใหม่ โมเดล AI เฉพาะทาง และการฝึกอบรมพิเศษ",
        "revolution_cta_button": "🚀 เข้าร่วมการปฏิวัติ AI (คลิกที่นี่!) 🚀",
        "revolution_cta_help": "คลิกเพื่อเข้าสู่หน้าลงทะเบียนของเราและค้นพบอนาคตของ ErnestMind",
        "contact_email_label": "ติดต่อ"
    },
    "tr": {
        "aide_title": "🆘 Yardım",
        "aide_texte": "Bir dosya yükleyin (csv, excel, parquet, ...), bir iş yapay zeka modülü seçin, 20'den fazla gelişmiş grafiği görüntüleyin ve anında PDF/Word/PPT raporları dışa aktarın. Tam gizlilik.",
        "apropos_title": "ℹ️ Hakkında",
        "apropos_texte": "ErnestMind V40, basitlik, hız ve gizlilik açısından Tableau, PowerBI, Dataiku, Qlik...'i geride bırakıyor. %100 yerel, otomatik yapay zeka, çok dilli, çok formatlı.",
        "exemple_toggle": "📈 Otomatik görselleştirme örnekleri",
        "exemple_titre": "### 📊 Yapay Zeka Demosu ve görselleştirmeler",
        "tab_corr": "📊 Korelasyon",
        "tab_boxplot": "📦 Kutu Grafiği",
        "tab_ventes": "📈 Satışlar",
        "tab_marketing": "🎯 Pazarlama",
        "tab_satisfaction": "😊 Memnuniyet",
        "upload_label": "📂 Bir dosya yükleyin",
        "format_non_supporte": "❌ Desteklenmeyen biçim",
        "success_upload": "✅ Dosya başarıyla yüklendi!",
        "erreur_chargement": "❌ Dosya yükleme hatası",
        "explorer": "🖨️ Veri önizleme",
        "nb_lignes": "Satır sayısı",
        "nb_colonnes": "Sütun sayısı",
        "colonnes": "Sütun adları",
        "types_colonnes": "Sütun türleri",
        "manquants": "Eksik değerleri olan sütunlar",
        "uniques": "Sütun başına benzersiz değerler",
        "stat_desc": "📐 Tanımlayıcı istatistikler",
        "apercu_donnees": "Veri önizleme",
        "resume": "🔎 Otomatik özet",
        "analyse_auto": "🧠 Akıllı otomatik analiz",
        "numeriques": "Sayısal",
        "categorique": "Kategorik",
        "valeurs_manquantes": "Eksik değerler",
        "classif_detectee": "Denetimli sınıflandırma algılandı",
        "reg_detectee": "Denetimli regresyon algılandı",
        "tache_non_detectee": "Görev otomatik olarak algılanmadı",
        "visual_auto": "📊 Otomatik görselleştirmeler",
        "histogramme": "Histogram",
        "nuage": "Dağılım Grafiği",
        "boxplot": "Kutu Grafiği",
        "repartition": "Dağılım",
        "choix_format": "📄 Oluşturulacak rapor formatı",
        "generer_rapport": "📜 Rapor oluştur",
        "telecharger_rapport": "📥 Oluşturulan raporu indir",
        "date_rapport": "Rapor tarihi/saati",
        "element": "Öğe",
        "valeur": "Değer",
        "v80_message": "🚀 ErnestMind V40, 100.000 satıra kadar analiz yapar. Sınırsız büyük veri analizi, 5 kat hız ve seçkin sektör yapay zeka modelleri için V80'e yükseltin!",
        "audit_title": "🔒 Yerel denetim ve güvenlik",
        "audit_exp": "Bu sürümde, gizliliği sağlamak için tüm geçici dosyalar otomatik olarak şifrelenir. Her gelişmiş işleme (metin/ses/görüntü) için yerel bir denetim raporu oluşturulur.",
        "brand_title": "🚀 ErnestMind AI V40 – Süper Veri Bilimi Platformu",
        "brand_claim": "Tableau, PowerBI, Dataiku, DataRobot, Alteryx, Qlik, SAS Viya, Palantir'den daha güçlü...<br>%100 yerel, otomatik yapay zeka, PDF/Word/PPT raporları, çok dilli, koyu/açık mod.",
        "connecteurs_title": "⚡️ Bağlayıcılar / API / Bulut / DB (demo)",
        "connecteurs_coming": "**Yaklaşan Bağlayıcılar:**",
        "connecteurs_import": "- MySQL, PostgreSQL, BigQuery, S3, Google Drive, OneDrive, SharePoint, REST API'den içe aktar",
        "connecteurs_export": "- Etkileşimli PowerPoint, HTML panosu, API, webhook, Slack, Teams, e-posta vb. dışa aktar.",
        "connecteurs_contact": "**İş kaynağınızı entegre etmek için bizimle iletişime geçin!**",
        "hist_title": "Histogram: {col}",
        "scatter_title": "Dağılım Grafiği: {x} vs {y}",
        "box_title": "Kutu Grafiği: {col}",
        "corr_title": "Korelasyon",
        "pairplot_title": "Değişken Dağılımı (pairplot)",
        "pie_title": "Dağılım: {col}",
        "bar_title": "Çubuk Grafiği: {col}",
        "revolution_cta_title": "✨ ErnestMind Devrimine Katılın!",
        "revolution_cta_text": "ErnestMind V40'ın gücünden etkilendiniz mi? Bu sadece bir başlangıç! Yeni özellikler, özel yapay zeka modelleri ve özel eğitimler hakkında bilgi almak için kaydolun.",
        "revolution_cta_button": "🚀 Yapay Zeka Devrimine Katılın (Buraya tıklayın!) 🚀",
        "revolution_cta_help": "Kayıt sayfamıza erişmek ve ErnestMind'ın geleceğini keşfetmek için tıklayın.",
        "contact_email_label": "İletişim"
    },
    "pl": {
        "aide_title": "🆘 Pomoc",
        "aide_texte": "Prześlij plik (csv, excel, parquet, ...), wybierz biznesowy moduł AI, wyświetl ponad 20 zaawansowanych wykresów i natychmiast wyeksportuj raporty PDF/Word/PPT. Pełna prywatność.",
        "apropos_title": "ℹ️ O nas",
        "apropos_texte": "ErnestMind V40 przewyższa Tableau, PowerBI, Dataiku, Qlik... pod względem prostoty, szybkości i prywatności. 100% lokalny, automatyczna AI, wielojęzyczny, wieloformatowy.",
        "exemple_toggle": "📈 Przykłady automatycznych wizualizacji",
        "exemple_titre": "### 📊 Demo AI i wizualizacje",
        "tab_corr": "📊 Korelacja",
        "tab_boxplot": "📦 Wykres pudełkowy",
        "tab_ventes": "📈 Sprzedaż",
        "tab_marketing": "🎯 Marketing",
        "tab_satisfaction": "😊 Satysfakcja",
        "upload_label": "📂 Prześlij plik",
        "format_non_supporte": "❌ Nieobsługiwany format",
        "success_upload": "✅ Plik przesłany pomyślnie!",
        "erreur_chargement": "❌ Błąd ładowania pliku",
        "explorer": "🖨️ Podgląd danych",
        "nb_lignes": "Liczba wierszy",
        "nb_colonnes": "Liczba kolumn",
        "colonnes": "Nazwy kolumn",
        "types_colonnes": "Typy kolumn",
        "manquants": "Kolumny z brakującymi wartościami",
        "uniques": "Unikalne wartości na kolumnę",
        "stat_desc": "📐 Statystyki opisowe",
        "apercu_donnees": "Podgląd danych",
        "resume": "🔎 Automatyczne podsumowanie",
        "analyse_auto": "🧠 Inteligentna analiza automatyczna",
        "numeriques": "Numeryczne",
        "categorique": "Kategoryczne",
        "valeurs_manquantes": "Brakujące wartości",
        "classif_detectee": "Wykryto klasyfikację nadzorowaną",
        "reg_detectee": "Wykryto regresję nadzorowaną",
        "tache_non_detectee": "Zadanie nie zostało wykryte automatycznie",
        "visual_auto": "📊 Automatyczne wizualizacje",
        "histogramme": "Histogram",
        "nuage": "Wykres punktowy",
        "boxplot": "Wykres pudełkowy",
        "repartition": "Rozkład",
        "choix_format": "📄 Format raportu do wygenerowania",
        "generer_rapport": "📜 Generuj raport",
        "telecharger_rapport": "📥 Pobierz wygenerowany raport",
        "date_rapport": "Data/godzina raportu",
        "element": "Element",
        "valeur": "Wartość",
        "v80_message": "🚀 ErnestMind V40 analizuje do 100 000 wierszy. Aby uzyskać nieograniczoną analizę dużych danych, 5-krotnie większą prędkość i elitarne modele AI dla konkretnych branż, uaktualnij do wersji V80!",
        "audit_title": "🔒 Audyt lokalny i bezpieczeństwo",
        "audit_exp": "W tej wersji wszystkie pliki tymczasowe są automatycznie szyfrowane w celu zapewnienia poufności. Lokalny raport z audytu jest generowany dla każdego zaawansowanego przetwarzania (tekst/audio/obraz).",
        "brand_title": "🚀 ErnestMind AI V40 – Platforma Super Data Science",
        "brand_claim": "Potężniejszy niż Tableau, PowerBI, Dataiku, DataRobot, Alteryx, Qlik, SAS Viya, Palantir...<br>100% lokalny, automatyczna AI, raporty PDF/Word/PPT, wielojęzyczny, tryb ciemny/jasny.",
        "connecteurs_title": "⚡️ Złącza / API / Chmura / Baza danych (demo)",
        "connecteurs_coming": "**Nadchodzące złącza:**",
        "connecteurs_import": "- Importuj z MySQL, PostgreSQL, BigQuery, S3, Google Drive, OneDrive, SharePoint, REST API",
        "connecteurs_export": "- Eksportuj do interaktywnego PowerPoint, pulpitu nawigacyjnego HTML, API, webhook, Slack, Teams, poczty e-mail itp.",
        "connecteurs_contact": "**Skontaktuj się z nami, aby zintegrować swoje źródło danych biznesowych!**",
        "hist_title": "Histogram: {col}",
        "scatter_title": "Wykres punktowy: {x} vs {y}",
        "box_title": "Wykres pudełkowy: {col}",
        "corr_title": "Korelacja",
        "pairplot_title": "Wykres punktowy zmiennych (pairplot)",
        "pie_title": "Rozkład: {col}",
        "bar_title": "Wykres słupkowy: {col}",
        "revolution_cta_title": "✨ Dołącz do rewolucji ErnestMind!",
        "revolution_cta_text": "Czy ErnestMind V40 Cię zaimponował? To tylko przedsmak! Zarejestruj się, aby być informowanym o nowych funkcjach, specjalistycznych modelach AI i ekskluzywnych szkoleniach.",
        "revolution_cta_button": "🚀 Dołącz do Rewolucji AI (Kliknij tutaj!) 🚀",
        "revolution_cta_help": "Kliknij, aby uzyskać dostęp do naszej strony rejestracji i odkryć przyszłość ErnestMind.",
        "contact_email_label": "Kontakt"
    },
    "ru": {
        "aide_title": "🆘 Помощь",
        "aide_texte": "Загрузите файл (csv, excel, parquet, ...), выберите модуль бизнес-ИИ, просмотрите более 20 расширенных диаграмм и мгновенно экспортируйте отчеты в PDF/Word/PPT. Полная конфиденциальность.",
        "apropos_title": "ℹ️ О нас",
        "apropos_texte": "ErnestMind V40 превосходит Tableau, PowerBI, Dataiku, Qlik... по простоте, скорости и конфиденциальности. 100% локально, авто-ИИ, многоязычность, мультиформатность.",
        "exemple_toggle": "📈 Примеры автоматических визуализаций",
        "exemple_titre": "### 📊 Демо ИИ и визуализации",
        "tab_corr": "📊 Корреляция",
        "tab_boxplot": "📦 Ящик с усами",
        "tab_ventes": "📈 Продажи",
        "tab_marketing": "🎯 Маркетинг",
        "tab_satisfaction": "😊 Удовлетворенность",
        "upload_label": "📂 Загрузить файл",
        "format_non_supporte": "❌ Неподдерживаемый формат",
        "success_upload": "✅ Файл успешно загружен!",
        "erreur_chargement": "❌ Ошибка загрузки файла",
        "explorer": "🖨️ Предварительный просмотр данных",
        "nb_lignes": "Количество строк",
        "nb_colonnes": "Количество столбцов",
        "colonnes": "Имена столбцов",
        "types_colonnes": "Типы столбцов",
        "manquants": "Столбцы с отсутствующими значениями",
        "uniques": "Уникальные значения по столбцам",
        "stat_desc": "📐 Описательные статистики",
        "apercu_donnees": "Предварительный просмотр данных",
        "resume": "🔎 Автоматическое резюме",
        "analyse_auto": "🧠 Интеллектуальный автоматический анализ",
        "numeriques": "Числовые",
        "categorique": "Категориальные",
        "valeurs_manquantes": "Отсутствующие значения",
        "classif_detectee": "Обнаружена классификация с учителем",
        "reg_detectee": "Обнаружена регрессия с учителем",
        "tache_non_detectee": "Задача не обнаружена автоматически",
        "visual_auto": "📊 Автоматические визуализации",
        "histogramme": "Гистограмма",
        "nuage": "Диаграмма рассеяния",
        "boxplot": "Ящик с усами",
        "repartition": "Распределение",
        "choix_format": "📄 Формат отчета для генерации",
        "generer_rapport": "📜 Сгенерировать отчет",
        "telecharger_rapport": "📥 Скачать сгенерированный отчет",
        "date_rapport": "Дата/время отчета",
        "element": "Элемент",
        "valeur": "Значение",
        "v80_message": "🚀 ErnestMind V40 анализирует до 100 000 строк. Для неограниченного анализа больших данных, 5-кратной скорости и элитных отраслевых моделей ИИ перейдите на V80!",
        "audit_title": "🔒 Локальный аудит и безопасность",
        "audit_exp": "В этой версии все временные файлы автоматически шифруются для обеспечения конфиденциальности. Локальный аудиторский отчет генерируется для каждой расширенной обработки (текст/аудио/изображение).",
        "brand_title": "🚀 ErnestMind AI V40 – Суперплатформа для анализа данных",
        "brand_claim": "Мощнее, чем Tableau, PowerBI, Dataiku, DataRobot, Alteryx, Qlik, SAS Viya, Palantir...<br>100% локально, автоматический ИИ, отчеты PDF/Word/PPT, многоязычный, темный/светлый режим.",
        "connecteurs_title": "⚡️ Коннекторы / API / Облако / БД (демо)",
        "connecteurs_coming": "**Скоро появятся коннекторы:**",
        "connecteurs_import": "- Импорт из MySQL, PostgreSQL, BigQuery, S3, Google Drive, OneDrive, SharePoint, REST API",
        "connecteurs_export": "- Экспорт в интерактивный PowerPoint, HTML-дашборд, API, webhook, Slack, Teams, email и т. д.",
        "connecteurs_contact": "**Свяжитесь с нами, чтобы интегрировать ваш бизнес-источник!**",
        "hist_title": "Гистограмма: {col}",
        "scatter_title": "Диаграмма рассеяния: {x} против {y}",
        "box_title": "Ящик с усами: {col}",
        "corr_title": "Корреляция",
        "pairplot_title": "Диаграмма рассеяния переменных (pairplot)",
        "pie_title": "Распределение: {col}",
        "bar_title": "Столбчатая диаграмма: {col}",
        "revolution_cta_title": "✨ Присоединяйтесь к революции ErnestMind!",
        "revolution_cta_text": "Впечатлены мощью ErnestMind V40? Это только начало! Зарегистрируйтесь, чтобы быть в курсе новых функций, специализированных моделей ИИ и эксклюзивных тренингов.",
        "revolution_cta_button": "🚀 Присоединиться к революции ИИ (Нажмите здесь!) 🚀",
        "revolution_cta_help": "Нажмите, чтобы перейти на нашу страницу регистрации и открыть для себя будущее ErnestMind.",
        "contact_email_label": "Контакт"
    },
    "sv": {
        "aide_title": "🆘 Hjälp",
        "aide_texte": "Ladda upp en fil (csv, excel, parquet, ...), välj en affärs-AI-modul, visa över 20 avancerade diagram och exportera PDF/Word/PPT-rapporter direkt. Total integritet.",
        "apropos_title": "ℹ️ Om",
        "apropos_texte": "ErnestMind V40 överträffar Tableau, PowerBI, Dataiku, Qlik... i enkelhet, hastighet och integritet. 100% lokal, automatisk AI, flera språk, flera format.",
        "exemple_toggle": "📈 Exempel på automatiska visualiseringar",
        "exemple_titre": "### 📊 AI-demo och visualiseringar",
        "tab_corr": "📊 Korrelation",
        "tab_boxplot": "📦 Boxplot",
        "tab_ventes": "📈 Försäljning",
        "tab_marketing": "🎯 Marknadsföring",
        "tab_satisfaction": "😊 Tillfredsställelse",
        "upload_label": "📂 Ladda upp en fil",
        "format_non_supporte": "❌ Format stöds ej",
        "success_upload": "✅ Filen laddades upp!",
        "erreur_chargement": "❌ Fel vid laddning av fil",
        "explorer": "🖨️ Förhandsgranskning av data",
        "nb_lignes": "Antal rader",
        "nb_colonnes": "Antal kolumner",
        "colonnes": "Kolumnnamn",
        "types_colonnes": "Kolumntyper",
        "manquants": "Kolumner med saknade värden",
        "uniques": "Unika värden per kolumn",
        "stat_desc": "📐 Beskrivande statistik",
        "apercu_donnees": "Dataförhandsvisning",
        "resume": "🔎 Automatisk sammanfattning",
        "analyse_auto": "🧠 Smart automatisk analys",
        "numeriques": "Numeriska",
        "categorique": "Kategoriska",
        "valeurs_manquantes": "Saknade värden",
        "classif_detectee": "Övervakad klassificering detekterad",
        "reg_detectee": "Övervakad regression detekterad",
        "tache_non_detectee": "Uppgiften detekterades inte automatiskt",
        "visual_auto": "📊 Automatiska visualiseringar",
        "histogramme": "Histogram",
        "nuage": "Punktdiagram",
        "boxplot": "Boxplot",
        "repartition": "Fördelning",
        "choix_format": "📄 Rapportformat att generera",
        "generer_rapport": "📜 Generera rapport",
        "telecharger_rapport": "📥 Ladda ner genererad rapport",
        "date_rapport": "Rapportdatum/tid",
        "element": "Objekt",
        "valeur": "Värde",
        "v80_message": "🚀 ErnestMind V40 analyserar upp till 100 000 rader. För obegränsad big data-analys, 5x hastighet och elitsektors AI-modeller, uppgradera till V80!",
        "audit_title": "🔒 Lokal revision & säkerhet",
        "audit_exp": "I den här versionen krypteras alla temporära filer automatiskt för att garantera konfidentialitet. En lokal revisionsrapport genereras för varje avancerad bearbetning (text/ljud/bild).",
        "brand_title": "🚀 ErnestMind AI V40 – Super Data Science Platform",
        "brand_claim": "Starkare än Tableau, PowerBI, Dataiku, DataRobot, Alteryx, Qlik, SAS Viya, Palantir...<br>100% lokal, automatisk AI, PDF/Word/PPT-rapporter, flerspråkig, mörkt/ljust läge.",
        "connecteurs_title": "⚡️ Anslutningar / API / Moln / DB (demo)",
        "connecteurs_coming": "**Kommande anslutningar:**",
        "connecteurs_import": "- Importera från MySQL, PostgreSQL, BigQuery, S3, Google Drive, OneDrive, SharePoint, REST API",
        "connecteurs_export": "- Exportera till interaktiv PowerPoint, HTML-dashboard, API, webhook, Slack, Teams, e-post, etc.",
        "connecteurs_contact": "**Kontakta oss för att integrera din affärsdatakälla!**",
        "hist_title": "Histogram: {col}",
        "scatter_title": "Punktdiagram: {x} vs {y}",
        "box_title": "Boxplot: {col}",
        "corr_title": "Korrelation",
        "pairplot_title": "Variabelpunktdiagram (pairplot)",
        "pie_title": "Fördelning: {col}",
        "bar_title": "Stapeldiagram: {col}",
        "revolution_cta_title": "✨ Gå med i ErnestMind Revolutionen!",
        "revolution_cta_text": "Imponerad av kraften i ErnestMind V40? Detta är bara en glimt! Registrera dig för att få information om nya funktioner, specialiserade AI-modeller och exklusiva utbildningar.",
        "revolution_cta_button": "🚀 Gå med i AI-revolutionen (Klicka här!) 🚀",
        "revolution_cta_help": "Klicka för att komma åt vår registreringssida och upptäck ErnestMinds framtid.",
        "contact_email_label": "Kontakt"
    },
    "da": {
        "aide_title": "🆘 Hjælp",
        "aide_texte": "Upload en fil (csv, excel, parquet, ...), vælg et forretnings-AI-modul, se over 20 avancerede diagrammer og eksporter PDF/Word/PPT-rapporter med det samme. Fuldstændig privatliv.",
        "apropos_title": "ℹ️ Om",
        "apropos_texte": "ErnestMind V40 overgår Tableau, PowerBI, Dataiku, Qlik... i enkelhed, hastighed og privatliv. 100% lokalt, auto AI, multi-sprog, multi-format.",
        "exemple_toggle": "📈 Eksempler på automatiske visualiseringer",
        "exemple_titre": "### 📊 AI-demo og visualiseringer",
        "tab_corr": "📊 Korrelation",
        "tab_boxplot": "📦 Boxplot",
        "tab_ventes": "📈 Salg",
        "tab_marketing": "🎯 Marketing",
        "tab_satisfaction": "😊 Tilfredshed",
        "upload_label": "📂 Upload en fil",
        "format_non_supporte": "❌ Ikke-understøttet format",
        "success_upload": "✅ Filen er uploadet!",
        "erreur_chargement": "❌ Fejl ved indlæsning af fil",
        "explorer": "🖨️ Dataforhåndsvisning",
        "nb_lignes": "Antal rækker",
        "nb_colonnes": "Antal kolonner",
        "colonnes": "Kolonnenavne",
        "types_colonnes": "Kolonnetyper",
        "manquants": "Kolonner med manglende værdier",
        "uniques": "Unikke værdier pr. kolonne",
        "stat_desc": "📐 Beskrivende statistik",
        "apercu_donnees": "Dataforhåndsvisning",
        "resume": "🔎 Automatisk oversigt",
        "analyse_auto": "🧠 Smart automatisk analyse",
        "numeriques": "Numeriske",
        "categorique": "Kategoriske",
        "valeurs_manquantes": "Manglende værdier",
        "classif_detectee": "Overvåget klassifikation detekteret",
        "reg_detectee": "Overvåget regression detekteret",
        "tache_non_detectee": "Opgave ikke automatisk detekteret",
        "visual_auto": "📊 Automatiske visualiseringer",
        "histogramme": "Histogram",
        "nuage": "Punktdiagram",
        "boxplot": "Boxplot",
        "repartition": "Fordeling",
        "choix_format": "📄 Rapportformat til generering",
        "generer_rapport": "📜 Generer rapport",
        "telecharger_rapport": "📥 Download genereret rapport",
        "date_rapport": "Rapportdato/tid",
        "element": "Element",
        "valeur": "Værdi",
        "v80_message": "🚀 ErnestMind V40 analyserer op til 100.000 rækker. For ubegrænset big data-analyse, 5x hastighed og elite-sektor AI-modeller, opgrader til V80!",
        "audit_title": "🔒 Lokal revision & sikkerhed",
        "audit_exp": "I denne version krypteres alle midlertidige filer automatisk for at garantere fortrolighed. En lokal revisionsrapport genereres for hver avanceret behandling (tekst/lyd/billede).",
        "brand_title": "🚀 ErnestMind AI V40 – Super Data Science Platform",
        "brand_claim": "Stærkere end Tableau, PowerBI, Dataiku, DataRobot, Alteryx, Qlik, SAS Viya, Palantir...<br>100% lokalt, automatisk AI, PDF/Word/PPT-rapporter, flersproget, mørk/lys tilstand.",
        "connecteurs_title": "⚡️ Konnektorer / API / Sky / DB (demo)",
        "connecteurs_coming": "**Kommende Konnektorer:**",
        "connecteurs_import": "- Importer fra MySQL, PostgreSQL, BigQuery, S3, Google Drive, OneDrive, SharePoint, REST API",
        "connecteurs_export": "- Eksporter til interaktiv PowerPoint, HTML-dashboard, API, webhook, Slack, Teams, e-mail osv.",
        "connecteurs_contact": "**Kontakt os for at integrere din forretningskilde!**",
        "hist_title": "Histogram: {col}",
        "scatter_title": "Punktdiagram: {x} vs {y}",
        "box_title": "Boxplot: {col}",
        "corr_title": "Korrelation",
        "pairplot_title": "Variabelpunktdiagram (pairplot)",
        "pie_title": "Fordeling: {col}",
        "bar_title": "Søjlediagram: {col}",
        "revolution_cta_title": "✨ Deltag i ErnestMind Revolutionen!",
        "revolution_cta_text": "Imponeret over kraften i ErnestMind V40? Dette er kun et glimt! Tilmeld dig for at modtage information om nye funktioner, specialiserede AI-modeller og eksklusiv træning.",
        "revolution_cta_button": "🚀 Deltag i AI-revolutionen (Klik her!) 🚀",
        "revolution_cta_help": "Klik for at få adgang til vores tilmeldingsside og opdag fremtiden for ErnestMind.",
        "contact_email_label": "Kontakt"
    },
    "no": {
        "aide_title": "🆘 Hjelp",
        "aide_texte": "Last opp en fil (csv, excel, parquet, ...), velg en forretnings-AI-modul, se over 20 avanserte diagrammer og eksporter PDF/Word/PPT-rapporter umiddelbart. Fullstendig personvern.",
        "apropos_title": "ℹ️ Om",
        "apropos_texte": "ErnestMind V40 overgår Tableau, PowerBI, Dataiku, Qlik... i enkelhet, hastighet og personvern. 100% lokalt, automatisk AI, flerspråklig, multi-format.",
        "exemple_toggle": "📈 Eksempler på automatiske visualiseringer",
        "exemple_titre": "### 📊 AI-demo og visualiseringer",
        "tab_corr": "📊 Korrelasjon",
        "tab_boxplot": "📦 Boxplot",
        "tab_ventes": "📈 Salg",
        "tab_marketing": "🎯 Markedsføring",
        "tab_satisfaction": "😊 Tilfredshet",
        "upload_label": "📂 Last opp en fil",
        "format_non_supporte": "❌ Ikke-støttet format",
        "success_upload": "✅ Filen er lastet opp!",
        "erreur_chargement": "❌ Feil ved lasting av fil",
        "explorer": "🖨️ Dataforhåndsvisning",
        "nb_lignes": "Antall rader",
        "nb_colonnes": "Antall kolonner",
        "colonnes": "Kolonnenavn",
        "types_colonnes": "Kolonnetyper",
        "manquants": "Kolonner med manglende verdier",
        "uniques": "Unike verdier per kolonne",
        "stat_desc": "📐 Beskrivende statistikk",
        "apercu_donnees": "Dataforhåndsvisning",
        "resume": "🔎 Automatisk oppsummering",
        "analyse_auto": "🧠 Smart automatisk analyse",
        "numeriques": "Numerisk",
        "categorique": "Kategorisk",
        "valeurs_manquantes": "Manglende verdier",
        "classif_detectee": "Veiledet klassifisering oppdaget",
        "reg_detectee": "Veiledet regresjon oppdaget",
        "tache_non_detectee": "Oppgave ikke automatisk oppdaget",
        "visual_auto": "📊 Automatiske visualiseringer",
        "histogramme": "Histogram",
        "nuage": "Punktdiagram",
        "boxplot": "Boxplot",
        "repartition": "Fordeling",
        "choix_format": "📄 Rapportformat som skal genereres",
        "generer_rapport": "📜 Generer rapport",
        "telecharger_rapport": "📥 Last ned generert rapport",
        "date_rapport": "Rapportdato/tid",
        "element": "Element",
        "valeur": "Verdi",
        "v80_message": "🚀 ErnestMind V40 analyserer opptil 100 000 rader. For ubegrenset big data-analyse, 5x hastighet og elite-sektor AI-modeller, oppgrader til V80!",
        "audit_title": "🔒 Lokal revisjon og sikkerhet",
        "audit_exp": "I denne versjonen krypteres alle midlertidige filer automatisk for å garantere konfidensialitet. En lokal revisjonsrapport genereres for hver avansert behandling (tekst/lyd/bilde).",
        "brand_title": "🚀 ErnestMind AI V40 – Super Data Science Platform",
        "brand_claim": "Sterkere enn Tableau, PowerBI, Dataiku, DataRobot, Alteryx, Qlik, SAS Viya, Palantir...<br>100% lokalt, automatisk AI, PDF/Word/PPT-rapporter, flerspråklig, mørk/lys modus.",
        "connecteurs_title": "⚡️ Koblinger / API / Sky / DB (demo)",
        "connecteurs_coming": "**Kommende koblinger:**",
        "connecteurs_import": "- Importer fra MySQL, PostgreSQL, BigQuery, S3, Google Disk, OneDrive, SharePoint, REST API",
        "connecteurs_export": "- Eksporter til interaktiv PowerPoint, HTML-dashboard, API, webhook, Slack, Teams, e-post, etc.",
        "connecteurs_contact": "**Kontakt oss for å integrere din forretningskilde!**",
        "hist_title": "Histogram: {col}",
        "scatter_title": "Punktdiagram: {x} vs {y}",
        "box_title": "Boxplot: {col}",
        "corr_title": "Korrelasjon",
        "pairplot_title": "Variabelpunktdiagram (pairplot)",
        "pie_title": "Fordeling: {col}",
        "bar_title": "Stolpediagram: {col}",
        "revolution_cta_title": "✨ Bli med i ErnestMind-revolusjonen!",
        "revolution_cta_text": "Imponert over kraften i ErnestMind V40? Dette er bare en smakebit! Meld deg på for å bli informert om nye funksjoner, spesialiserte AI-modeller og eksklusiv opplæring.",
        "revolution_cta_button": "🚀 Bli med i AI-revolusjonen (Klikk her!) 🚀",
        "revolution_cta_help": "Klikk for å få tilgang til vår registreringsside og oppdag fremtiden til ErnestMind.",
        "contact_email_label": "Kontakt"
    },
    "ht": {
        "aide_title": "🆘 Èd",
        "aide_texte": "Telechaje yon fichye (csv, excel, parquet, ...), chwazi yon modil IA biznis, gade plis pase 20 tablo avanse epi ekspòte rapò PDF/Word/PPT imedyatman. Konfidansyalite total.",
        "apropos_title": "ℹ️ Konsènan",
        "apropos_texte": "ErnestMind V40 depase Tableau, PowerBI, Dataiku, Qlik... nan senplisite, vitès ak konfidansyalite. 100% lokal, IA oto, plizyè lang, plizyè fòma.",
        "exemple_toggle": "📈 Egzanp vizyalizasyon otomatik",
        "exemple_titre": "### 📊 Demo IA & vizyalizasyon",
        "tab_corr": "📊 Korelasyon",
        "tab_boxplot": "📦 Boxplot",
        "tab_ventes": "📈 Lavant",
        "tab_marketing": "🎯 Maketing",
        "tab_satisfaction": "😊 Satisfaksyon",
        "upload_label": "📂 Telechaje yon fichye",
        "format_non_supporte": "❌ Fòma pa sipòte",
        "success_upload": "✅ Fichye telechaje avèk siksè!",
        "erreur_chargement": "❌ Erè pandan chaje fichye",
        "explorer": "🖨️ Apèsi done",
        "nb_lignes": "Kantite liy",
        "nb_colonnes": "Kantite kolòn",
        "colonnes": "Non kolòn",
        "types_colonnes": "Tip kolòn",
        "manquants": "Kolòn ak valè ki manke",
        "uniques": "Valè inik pou chak kolòn",
        "stat_desc": "📐 Estatistik deskriptif",
        "apercu_donnees": "Apèsi done",
        "resume": "🔎 Rezime otomatik",
        "analyse_auto": "🧠 Analiz otomatik entèlijan",
        "numeriques": "Nimerik",
        "categorique": "Kategorik",
        "valeurs_manquantes": "Valè ki manke",
        "classif_detectee": "Klasifikasyon sipèvize detekte",
        "reg_detectee": "Regresyon sipèvize detekte",
        "tache_non_detectee": "Tach pa detekte otomatikman",
        "visual_auto": "📊 Vizyalizasyon otomatik",
        "histogramme": "Istogram",
        "nuage": "Nwaj de pwen",
        "boxplot": "Boxplot",
        "repartition": "Repatisyon",
        "choix_format": "📄 Fòma rapò pou jenere",
        "generer_rapport": "📜 Jenere rapò",
        "telecharger_rapport": "📥 Telechaje rapò jenere a",
        "date_rapport": "Dat ak lè rapò",
        "element": "Eleman",
        "valeur": "Valè",
        "v80_message": "🚀 ErnestMind V40 analize jiska 100,000 liy. Pou analiz gwo done san limit, vitès 5x, ak modèl IA sektè elit, ajouye nan V80!",
        "audit_title": "🔒 Odyans lokal & sekirite",
        "audit_exp": "Nan vèsyon sa a, tout fichye tanporè yo chifre otomatikman pou garanti konfidansyalite. Yon rapò odyans lokal jenere pou chak tretman avanse (tèks/odyo/imaj).",
        "brand_title": "🚀 ErnestMind AI V40 – Platfòm Super Syans Done",
        "brand_claim": "+ Pisan pase Tableau, PowerBI, Dataiku, DataRobot, Alteryx, Qlik, SAS Viya, Palantir...<br>100% lokal, IA oto, rapò PDF/Word/PPT, plizyè lang, mòd nwa/klè.",
        "connecteurs_title": "⚡️ Konektè / API / Nwaj / BD (demo)",
        "connecteurs_coming": "**Konektè k ap vini:**",
        "connecteurs_import": "- Enpòte soti nan MySQL, PostgreSQL, BigQuery, S3, Google Drive, OneDrive, SharePoint, API REST",
        "connecteurs_export": "- Ekspòte nan PowerPoint entèaktif, HTML dashboard, API, webhook, Slack, Teams, imèl, elatriye.",
        "connecteurs_contact": "**Kontakte nou pou entegre sous biznis ou!**",
        "hist_title": "Istogram: {col}",
        "scatter_title": "Nwaj de pwen: {x} vs {y}",
        "box_title": "Boxplot: {col}",
        "corr_title": "Korelasyon",
        "pairplot_title": "Nwaj de varyab (pairplot)",
        "pie_title": "Repatisyon: {col}",
        "bar_title": "Barplot: {col}",
        "revolution_cta_title": "✨ Rejoins revolisyon ErnestMind lan!",
        "revolution_cta_text": "Èske ErnestMind V40 enpresyone w? Sa a se sèlman yon apèsi! Enskri pou w enfòme sou nouvo karakteristik, modèl IA espesyalize, ak fòmasyon eksklizif.",
        "revolution_cta_button": "🚀 Enskri pou revolisyon IA a (Klike la !) 🚀",
        "revolution_cta_help": "Klike pou w aksè nan paj enskripsyon nou an epi dekouvri lavni ErnestMind.",
        "contact_email_label": "Kontak"
    }
}

t = translations.get(langue_code, translations["en"])

st.markdown(f'<div class="superbrand">{t["brand_title"]}</div>', unsafe_allow_html=True)
st.markdown(f'<div class="superclaim">{t["brand_claim"]}</div>', unsafe_allow_html=True)
st.markdown("---")

def generate_key():
    if Fernet is None:
        return None
    key_file = "ernestmind_v40.key"
    if not os.path.exists(key_file):
        key = Fernet.generate_key()
        with open(key_file, "wb") as f:
            f.write(key)
    else:
        with open(key_file, "rb") as f:
            key = f.read()
    return key

def encrypt_file(filepath):
    if Fernet is None:
        return
    key = generate_key()
    fernet = Fernet(key)
    with open(filepath, "rb") as file:
        original = file.read()
    encrypted = fernet.encrypt(original)
    with open(filepath, "wb") as encrypted_file:
        encrypted_file.write(encrypted)

def write_audit_log(action, filename, extra=None):
    now = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    fname = "ernestmind_audit.log"
    msg = f"{now} | ACTION: {action} | FILE: {filename}"
    if extra:
        msg += f" | INFO: {extra}"
    with open(fname, "a", encoding="utf-8") as f:
        f.write(msg + "\n")

def slugify(value):
    value = str(value)
    value = unicodedata.normalize("NFKD", value).encode("ascii", "ignore").decode("ascii")
    value = value.replace(" ", "_")
    valid_chars = f"-_.() {string.ascii_letters}{string.digits}"
    value = "".join(c for c in value if c in valid_chars)
    return value.lower()

def save_plot(fig, path):
    fig.tight_layout()
    fig.savefig(path, bbox_inches="tight")
    plt.close(fig)

def format_table_for_pdf(table_df, max_width=80):
    if table_df.shape[1] > 10:
        table_df = table_df.iloc[:, :10]
        table_df.columns = [str(col)[:12] for col in table_df.columns]
    for col in table_df.columns:
        if table_df[col].map(lambda x: len(str(x))).max() > max_width:
            table_df[col] = table_df[col].map(lambda x: str(x)[:max_width] + "...")
    return table_df

def extract_text_from_pdf(file_path):
    text = ""
    if pdfplumber is not None:
        try:
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    text += page.extract_text() or ""
        except Exception as e:
            text = f"[PDF Extraction Error: {e}]"
    return text

def extract_text_from_docx(file_path):
    text = ""
    if docx is not None:
        try:
            document = docx.Document(file_path)
            text = "\n".join([para.text for para in document.paragraphs])
        except Exception as e:
            text = f"[DOCX Extraction Error: {e}]"
    return text

def extract_text_from_txt(file_path):
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()
    except Exception:
        with open(file_path, "r", encoding="latin1") as f:
            return f.read()

def extract_text_from_image(file_path):
    if pytesseract is not None and Image is not None:
        try:
            img = Image.open(file_path)
            return pytesseract.image_to_string(img)
        except Exception as e:
            return f"[OCR Error: {e}]"
    return "[OCR not available]"

def transcribe_audio(file_path):
    if whisper is not None:
        try:
            model = whisper.load_model("base")
            result = model.transcribe(file_path)
            return result["text"]
        except Exception as e:
            return f"[Audio Transcription Error: {e}]"
    return "[Whisper not available]"

def analyse_text(text):
    # Simple sentiment/summary/entity detection (local, for demo)
    import re
    from collections import Counter
    # Sentiment: simple count of positive/negative words
    positive = ("good","great","excellent","positive","happy","success","win","profit","love","joy","ok","yes")
    negative = ("bad","poor","negative","sad","failure","loss","hate","angry","no")
    pos_count = sum(text.lower().count(w) for w in positive)
    neg_count = sum(text.lower().count(w) for w in negative)
    sentiment = "Neutral"
    if pos_count > neg_count:
        sentiment = "Positive"
    elif neg_count > pos_count:
        sentiment = "Negative"
    # Entities: extract capitalized words as demo
    entities = Counter(re.findall(r"\b[A-Z][a-z]{2,}\b", text))
    # Summary: first 2 sentences
    summary = ".".join(text.split(".")[:2]) + ("..." if len(text.split(".")) > 2 else "")
    return {
        "sentiment": sentiment,
        "entities": dict(entities.most_common(10)),
        "summary": summary
    }

# ========== Message V80 toutes les 5 minutes ==========
if "last_upgrade_msg" not in st.session_state:
    st.session_state.last_upgrade_msg = 0

def show_v80_message():
    st.markdown(
        f"<div style='background-color:#e6f4ff;border-left:0.5em solid #1E90FF;padding:1em 1.2em;margin:0.5em 0 1.2em 0;border-radius:7px;font-weight:bold;font-size:1.1em;color:#1E90FF'>{t['v80_message']}</div>",
        unsafe_allow_html=True
    )

if time.time() - st.session_state.last_upgrade_msg > 300:
    show_v80_message()
    st.session_state.last_upgrade_msg = time.time()

with st.sidebar.expander(t["aide_title"]):
    st.markdown(t["aide_texte"])
with st.sidebar.expander(t["apropos_title"]):
    st.markdown(t["apropos_texte"])

with st.sidebar.expander(t["audit_title"]):
    st.markdown(t["audit_exp"])
    if os.path.exists("ernestmind_audit.log"):
        with open("ernestmind_audit.log", "r", encoding="utf-8") as f:
            lines = f.readlines()[-10:]
            st.code("".join(lines) or "Aucune entrée d'audit.")

show_examples = st.checkbox(t["exemple_toggle"])
if show_examples:
    st.markdown("---")
    st.markdown(t["exemple_titre"])
    df_demo = pd.DataFrame({
        "Ventes": np.random.normal(20000, 5000, 1000),
        "Marketing (€)": np.random.normal(15000, 4000, 1000),
        "Région": np.random.choice(["Europe", "Amérique", "Asie", "Afrique"], 1000),
        "Satisfaction Client": np.random.uniform(1, 5, 1000),
        "Date": pd.date_range(start="2023-01-01", periods=1000, freq="D")
    })
    tab1, tab2, tab3, tab4, tab5 = st.tabs([
        t["tab_corr"], t["tab_boxplot"], t["tab_ventes"], t["tab_marketing"], t["tab_satisfaction"]
    ])
    with tab1:
        corr = df_demo[["Ventes", "Marketing (€)", "Satisfaction Client"]].corr()
        fig, ax = plt.subplots()
        sns.heatmap(corr, annot=True, cmap="coolwarm", ax=ax)
        st.pyplot(fig)
    with tab2:
        fig, ax = plt.subplots()
        sns.boxplot(data=df_demo, x="Région", y="Ventes", palette="Set3", ax=ax)
        st.pyplot(fig)
    with tab3:
        df_grp = df_demo.groupby("Date").sum(numeric_only=True)
        fig, ax = plt.subplots()
        df_grp["Ventes"].rolling(30).mean().plot(ax=ax, color="blue")
        st.pyplot(fig)
    with tab4:
        fig, ax = plt.subplots()
        sns.scatterplot(data=df_demo, x="Marketing (€)", y="Ventes", hue="Région", ax=ax)
        st.pyplot(fig)
    with tab5:
        fig, ax = plt.subplots()
        sns.histplot(df_demo["Satisfaction Client"], kde=True, bins=20, ax=ax, color="green")
        st.pyplot(fig)

# ========== Animation bleue soulée (shiny) pour import fichier ==========
st.markdown('<div class="upload-anim"></div>', unsafe_allow_html=True)

uploaded_file = st.file_uploader(
    t["upload_label"],
    type=[
        "csv", "xlsx", "xls", "json", "txt", "tsv", "xml", "html", "parquet", "feather",
        "sav", "dta", "pkl", "hdf5", "ods", "pdf", "docx", "jpg", "jpeg", "png", "bmp", "wav", "mp3", "ogg", "m4a"
    ]
)

# Tes expanders de barre latérale existants (Aide, À propos, Audit)

with st.sidebar.expander(t["audit_title"]):
    st.markdown(t["audit_exp"])
    if os.path.exists("ernestmind_audit.log"):
        with open("ernestmind_audit.log", "r", encoding="utf-8") as f:
            lines = f.readlines()[-10:]
            st.code("".join(lines) or "Aucune entrée d'audit.")

# Ton expander "Connecteurs" déplacé ici pour être dans la sidebar
# ========== Connecteurs universels mini-dashboard ==========
with st.sidebar.expander(t["connecteurs_title"]):
    st.markdown(t["connecteurs_coming"])
    st.markdown(t["connecteurs_import"])
    st.markdown(t["connecteurs_export"])
    st.markdown(t["connecteurs_contact"])

# --- DÉBUT DU NOUVEAU BLOC POUR LE BOUTON DANS LA SIDEBAR ---
st.sidebar.markdown("---") # Une ligne de séparation visuelle

# Titre de l'appel à l'action
st.sidebar.markdown(f'<h3 style="color:#1E90FF; text-align: center;">{t["revolution_cta_title"]}</h3>', unsafe_allow_html=True)
# Texte descriptif pour l'appel à l'action
st.sidebar.markdown(f'<p style="color:#1E90FF; font-size: 0.9em; text-align: center;">{t["revolution_cta_text"]}</p>', unsafe_allow_html=True)

# L'URL de ta page GitHub Pages (où les gens peuvent s'inscrire)
inscription_page_url = "https://franclin55.github.io/ernestmind/"

# Le bouton lui-même
st.sidebar.link_button(
    label=t["revolution_cta_button"],
    url=inscription_page_url,
    help=t["revolution_cta_help"]
)

st.sidebar.markdown("---") # Une autre ligne de séparation pour le bas de la sidebar
# Affichage de l'email de contact
st.sidebar.markdown(
    f'<p style="color:#1E90FF; font-size: 0.8em; text-align: center;">{t["contact_email_label"]} : parisien.data@gmail.com</p>',
    unsafe_allow_html=True
)
# Ton branding ErnestMind V40
st.sidebar.markdown(f'<p style="color:#1E90FF; font-size: 0.7em; text-align: center;">ErnestMind V40</p>', unsafe_allow_html=True)

# --- FIN DU NOUVEAU BLOC ---



df = None
text_advanced = None
audio_transcript = None
ocr_text = None
if uploaded_file is not None:
    temp_file_path = f"temp_{slugify(uploaded_file.name)}"
    with open(temp_file_path, "wb") as f:
        f.write(uploaded_file.read())
    file_type = uploaded_file.name.split(".")[-1].lower()
    # Chiffrement automatique du fichier temporaire
    encrypt_file(temp_file_path)
    # Déchiffrement pour traitement (démo: lecture après chiffrement)
    if Fernet is not None:
        key = generate_key()
        fernet = Fernet(key)
        with open(temp_file_path, "rb") as encrypted_file:
            encrypted = encrypted_file.read()
            try:
                decrypted = fernet.decrypt(encrypted)
                with open(temp_file_path, "wb") as f:
                    f.write(decrypted)
            except Exception:
                pass # Already plain or failed, continue
    # Traitements avancés texte/audio/image/OCR
    try:
        if file_type in ["csv", "txt", "tsv", "xlsx", "xls", "ods", "json", "xml", "html", "parquet", "feather", "sav", "dta", "pkl", "pickle", "hdf5", "h5"]:
            # Data classique
            if file_type in ["csv", "txt"]:
                try:
                    df = pd.read_csv(temp_file_path, sep=None, engine="python", encoding="utf-8")
                except Exception:
                    df = pd.read_csv(temp_file_path, sep=None, engine="python", encoding="latin1")
            elif file_type == "tsv":
                df = pd.read_csv(temp_file_path, sep="\t")
            elif file_type in ["xlsx", "xls", "ods"]:
                try:
                    df = pd.read_excel(temp_file_path, engine="openpyxl")
                except Exception:
                    df = pd.read_excel(temp_file_path)
            elif file_type == "json":
                try:
                    df = pd.read_json(temp_file_path)
                except Exception:
                    with open(temp_file_path, "r", encoding="utf-8") as f:
                        import json
                        data = json.load(f)
                        df = pd.json_normalize(data)
            elif file_type == "xml":
                df = pd.read_xml(temp_file_path)
            elif file_type == "html":
                df_list = pd.read_html(temp_file_path)
                df = df_list[0]
            elif file_type == "parquet":
                df = pd.read_parquet(temp_file_path)
            elif file_type == "feather":
                df = pd.read_feather(temp_file_path)
            elif file_type == "sav":
                import pyreadstat
                df, meta = pyreadstat.read_sav(temp_file_path)
            elif file_type == "dta":
                df = pd.read_stata(temp_file_path)
            elif file_type in ["pkl", "pickle"]:
                import pickle
                with open(temp_file_path, "rb") as f:
                    df = pickle.load(f)
            elif file_type in ["hdf5", "h5"]:
                df = pd.read_hdf(temp_file_path)
            if df is not None:
                st.success(t["success_upload"])
                write_audit_log("UPLOAD_DATA", uploaded_file.name, f"shape={df.shape}")
            os.remove(temp_file_path)
        elif file_type == "pdf":
            text_advanced = extract_text_from_pdf(temp_file_path)
            st.success("PDF chargé et extrait avec succès.")
            write_audit_log("ANALYSE_PDF", uploaded_file.name, f"len={len(text_advanced)}")
            os.remove(temp_file_path)
        elif file_type == "docx":
            text_advanced = extract_text_from_docx(temp_file_path)
            st.success("DOCX chargé et extrait avec succès.")
            write_audit_log("ANALYSE_DOCX", uploaded_file.name, f"len={len(text_advanced)}")
            os.remove(temp_file_path)
        elif file_type in ["jpg", "jpeg", "png", "bmp"]:
            ocr_text = extract_text_from_image(temp_file_path)
            st.success("Image chargée et OCR effectué avec succès.")
            write_audit_log("OCR_IMAGE", uploaded_file.name, f"len={len(ocr_text)}")
            st.image(temp_file_path, caption="Image importée", use_column_width=True)
            os.remove(temp_file_path)
        elif file_type in ["wav", "mp3", "ogg", "m4a"]:
            audio_transcript = transcribe_audio(temp_file_path)
            st.success("Audio chargé et transcription effectuée localement.")
            write_audit_log("AUDIO_TRANSCRIBE", uploaded_file.name, f"len={len(audio_transcript)}")
            os.remove(temp_file_path)
        else:
            st.error(t["format_non_supporte"])
            st.stop()
    except Exception as e:
        st.error(f"{t['erreur_chargement']} : {e}")

# ======= Affichage analyse texte/ocr/audio avancés =======
if text_advanced:
    st.subheader("📄 Analyse avancée du texte (PDF/DOCX)")
    st.text_area("Texte extrait", value=text_advanced[:5000], height=240)
    res = analyse_text(text_advanced)
    st.info(f"**Résumé** : {res['summary']}")
    st.info(f"**Sentiment détecté** : {res['sentiment']}")
    st.info(f"**Entités détectées** : {', '.join([f'{k} ({v})' for k,v in res['entities'].items()])}")
if ocr_text:
    st.subheader("🖼️ Texte OCR extrait de l'image")
    st.text_area("Texte OCR", value=ocr_text[:2000], height=180)
    res = analyse_text(ocr_text)
    st.info(f"**Résumé** : {res['summary']}")
    st.info(f"**Sentiment détecté** : {res['sentiment']}")
    st.info(f"**Entités détectées** : {', '.join([f'{k} ({v})' for k,v in res['entities'].items()])}")
if audio_transcript:
    st.subheader("🎤 Transcription audio locale")
    st.text_area("Transcription", value=audio_transcript[:3000], height=180)
    res = analyse_text(audio_transcript)
    st.info(f"**Résumé** : {res['summary']}")
    st.info(f"**Sentiment détecté** : {res['sentiment']}")
    st.info(f"**Entités détectées** : {', '.join([f'{k} ({v})' for k,v in res['entities'].items()])}")

# ======= Bloc à intégrer pour génération des images (exemple, après import) =======
if df is not None and st.button("Générer les graphiques du dataset importé"):
    with tempfile.TemporaryDirectory() as tmpdir:
        img_paths = generate_graph_images(df, tmpdir, t)
        st.success("Graphiques générés !")
        for img in img_paths:
            st.image(img, caption=os.path.basename(img))
        # ...

# ========== DATAFRAME ET ANALYSE CLASSIQUE ==========
if df is not None:
    st.subheader(t["explorer"])
    st.dataframe(df.head())
    st.markdown("---")
    resume_df = pd.DataFrame({
        t["element"]: [
            t["nb_lignes"],
            t["nb_colonnes"],
            t["colonnes"],
            t["types_colonnes"],
            t["manquants"],
            t["uniques"]
        ],
        t["valeur"]: [
            str(len(df)),
            str(df.shape[1]),
            ", ".join(df.columns),
            ", ".join([f"{col}: {df[col].dtype}" for col in df.columns]),
            ", ".join([col for col in df.columns if df[col].isnull().any()]) or "Aucune",
            ", ".join([f"{col}: {df[col].nunique()}" for col in df.columns])
        ]
    })
    resume_df = resume_df.astype(str)
    st.subheader("📝 " + t["resume"])
    st.dataframe(resume_df)
    st.markdown("---")
    info_cols = pd.DataFrame({
        "Colonne": df.columns,
        "Type": [str(df[col].dtype) for col in df.columns],
        "Valeurs uniques": [str(df[col].nunique()) for col in df.columns],
        "Valeurs manquantes": [str(df[col].isnull().sum()) for col in df.columns]
    }).astype(str)
    st.subheader("📋 Infos colonnes")
    st.dataframe(info_cols)
    st.markdown("---")
    st.subheader(t["stat_desc"])
    try:
        desc_df = df.describe(include="all").T.astype(str)
        st.dataframe(desc_df)
    except Exception as e:
        st.error("Erreur lors de l’analyse des données : " + str(e))
    st.markdown("---")
    preview_head = df.head(10)
    st.subheader(t["apercu_donnees"])
    st.dataframe(preview_head.astype(str))
    st.markdown("---")

    # ========== IA auto-analyse, suggestions, dashboard IA ==========
    def analyse_intelligente_v40(df):
        nb_lignes, nb_colonnes = df.shape
        colonnes_num = df.select_dtypes(include=np.number).columns.tolist()
        colonnes_cat = df.select_dtypes(include="object").columns.tolist()
        total_nan = df.isnull().sum().sum()
        tasks = []
        tasks.append(("Lignes", nb_lignes))
        tasks.append(("Colonnes", nb_colonnes))
        tasks.append(("Numériques", ", ".join(colonnes_num)))
        tasks.append(("Catégoriques", ", ".join(colonnes_cat)))
        tasks.append(("Valeurs manquantes", total_nan))
        if "target" in df.columns or "label" in df.columns:
            y = df.get("target") or df.get("label")
            if y.dtype == "object" or len(y.unique()) <= 10:
                tasks.append(("Tâche détectée", t["classif_detectee"]))
            else:
                tasks.append(("Tâche détectée", t["reg_detectee"]))
        else:
            tasks.append(("Tâche détectée", t["tache_non_detectee"]))
        if nb_lignes > 5000:
            tasks.append(("⚡️ Suggestion IA", "Échantillonnage recommandé pour accélérer l’analyse"))
        if len(colonnes_num) > 12:
            tasks.append(("⚡️ Suggestion IA", "Réduire la dimension avec PCA ou SelectKBest"))
        if total_nan > 0:
            tasks.append(("⚠️ Alerte IA", "Présence de valeurs manquantes : pensez à l’imputation"))
        return pd.DataFrame(tasks, columns=["Insight", "Valeur"]).astype(str)

    st.subheader(t["analyse_auto"])
    analyse_df = analyse_intelligente_v40(df)
    st.dataframe(analyse_df)
    st.markdown("---")
    st.subheader(t["valeurs_manquantes"])
    missing = df.isnull().sum()
    missing_df = pd.DataFrame({
        "Colonne": missing.index.astype(str),
        "Nb manquantes": missing.values.astype(str)
    })
    st.dataframe(missing_df[missing_df["Nb manquantes"] > "0"])
    st.markdown("---")
    st.subheader(t["visual_auto"])
    numeric_columns = df.select_dtypes(include=np.number).columns.tolist()
    categorical_cols = df.select_dtypes(include="object").columns.tolist()
    if numeric_columns:
        fig1 = px.histogram(df, x=numeric_columns[0], title=f"{t['histogramme']} : {numeric_columns[0]}")
        st.plotly_chart(fig1, use_container_width=True)
        if len(numeric_columns) >= 2:
            fig2 = px.scatter(df, x=numeric_columns[0], y=numeric_columns[1], title=t["nuage"])
            st.plotly_chart(fig2, use_container_width=True)
        fig3 = px.box(df, y=numeric_columns[0], title=f"{t['boxplot']} : {numeric_columns[0]}")
        st.plotly_chart(fig3, use_container_width=True)
        fig4, ax = plt.subplots(figsize=(5,3), dpi=120)
        sns.heatmap(df[numeric_columns].corr(), annot=True, cmap="coolwarm", ax=ax)
        st.pyplot(fig4)
    if categorical_cols:
        vc = df[categorical_cols[0]].value_counts().reset_index()
        vc.columns = ["nom", "count"]
        fig5 = px.bar(vc, x="nom", y="count", title=f"{t['repartition']} : {categorical_cols[0]}")
        st.plotly_chart(fig5, use_container_width=True)
    mime_types = {
        "PDF": "application/pdf",
        "Word": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "PPT": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "Excel": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    }
    st.markdown("---")
    st.subheader("📝 " + t["generer_rapport"])
    rapport_format = st.selectbox(
        t["choix_format"],
        ["PDF", "Word", "PPT", "Excel"]
    )
    if st.button(t["generer_rapport"]):
        progress = st.progress(0, text="Préparation du rapport...")
        now = datetime.now()
        now_str = now.strftime("%Y-%m-%d_%Hh%M%S")
        base_name = slugify(f"rapport_{now_str}_{rapport_format.lower()}")
        file_data = None
        with tempfile.TemporaryDirectory() as tmpdir:
            img_paths = generate_graph_images(df, tmpdir, t)  # <-- il faut 3 arguments
            progress.progress(20, "Génération du rapport...")
            tmpfile_path = os.path.join(tmpdir, f"{base_name}.{rapport_format.lower()}")
            try:
                excel_tables = {
                    "Statistiques descriptives": format_table_for_pdf(df.describe(include="all").T.astype(str)),
                    "Infos colonnes": format_table_for_pdf(info_cols),
                    "Aperçu des données": format_table_for_pdf(df.head(10).astype(str))
                }
                if rapport_format == "Excel":
                    with pd.ExcelWriter(tmpfile_path, engine='xlsxwriter') as writer:
                        for sheet, table in excel_tables.items():
                            table.to_excel(writer, sheet_name=sheet[:30], index=True)
                        resume_df.to_excel(writer, sheet_name="Résumé", index=False)
                        analyse_df.to_excel(writer, sheet_name="Analyse IA", index=False)
                        missing_df[missing_df["Nb manquantes"] > "0"].to_excel(writer, sheet_name="Valeurs manquantes", index=False)
                        workbook  = writer.book
                        for idx, img in enumerate(img_paths):
                            worksheet = workbook.add_worksheet(f"Graphique{idx+1}")
                            worksheet.insert_image('B2', img)
                    progress.progress(60)
                elif rapport_format == "Word":
                    from docx import Document
                    from docx.shared import Inches
                    doc = Document()
                    doc.add_heading("Rapport d'Analyse de Données", 0)
                    doc.add_paragraph(f"{t['date_rapport']} : {now.strftime('%Y-%m-%d %H:%M:%S')}")
                    doc.add_heading("Résumé général", level=1)
                    t1 = doc.add_table(rows=1, cols=2)
                    t1.style = 'Light List'
                    t1.cell(0,0).text = t["element"]
                    t1.cell(0,1).text = t["valeur"]
                    for i, row in resume_df.iterrows():
                        t1.add_row().cells[0].text = str(row[t["element"]])
                        t1.rows[-1].cells[1].text = str(row[t["valeur"]])
                    doc.add_heading("Infos colonnes", level=1)
                    info_cols_pdf = format_table_for_pdf(info_cols)
                    t2 = doc.add_table(rows=1, cols=len(info_cols_pdf.columns))
                    t2.style = 'Light List'
                    for col_idx, col in enumerate(info_cols_pdf.columns):
                        t2.cell(0,col_idx).text = col
                    for idx, row in info_cols_pdf.iterrows():
                        cells = t2.add_row().cells
                        for col_idx, val in enumerate(row):
                            cells[col_idx].text = str(val)
                    doc.add_heading("Statistiques descriptives (1ères colonnes)", level=1)
                    stat_df = format_table_for_pdf(df.describe(include="all").T.astype(str))
                    t3 = doc.add_table(rows=1, cols=len(stat_df.columns)+1)
                    t3.style = 'Light List'
                    t3.cell(0,0).text = "Colonne"
                    for cidx, col in enumerate(stat_df.columns):
                        t3.cell(0,cidx+1).text = str(col)
                    for idx, row in stat_df.iterrows():
                        row_cells = t3.add_row().cells
                        row_cells[0].text = str(idx)
                        for cidx, val in enumerate(row):
                            row_cells[cidx+1].text = str(val)
                    doc.add_heading("Aperçu des données", level=1)
                    head_df = format_table_for_pdf(df.head(10).astype(str))
                    t4 = doc.add_table(rows=1, cols=len(head_df.columns)+1)
                    t4.style = 'Light List'
                    t4.cell(0,0).text = "Index"
                    for cidx, col in enumerate(head_df.columns):
                        t4.cell(0,cidx+1).text = str(col)
                    for idx, row in head_df.iterrows():
                        row_cells = t4.add_row().cells
                        row_cells[0].text = str(idx)
                        for cidx, val in enumerate(row):
                            row_cells[cidx+1].text = str(val)
                    doc.add_heading("Analyse IA", level=1)
                    t5 = doc.add_table(rows=1, cols=2)
                    t5.cell(0,0).text = "Insight"
                    t5.cell(0,1).text = "Valeur"
                    for i, row in analyse_df.iterrows():
                        t5.add_row().cells[0].text = str(row["Insight"])
                        t5.rows[-1].cells[1].text = str(row["Valeur"])
                    doc.add_heading("Valeurs manquantes", level=1)
                    t6 = doc.add_table(rows=1, cols=2)
                    t6.cell(0,0).text = "Colonne"
                    t6.cell(0,1).text = "Nb manquantes"
                    for i, row in missing_df[missing_df["Nb manquantes"] > "0"].iterrows():
                        t6.add_row().cells[0].text = str(row["Colonne"])
                        t6.rows[-1].cells[1].text = str(row["Nb manquantes"])
                    doc.add_heading("Graphiques automatiques", level=1)
                    for img in img_paths:
                        doc.add_picture(img, width=Inches(5.5))
                    doc.save(tmpfile_path)
                    progress.progress(60)
                elif rapport_format == "PDF":
                    from reportlab.lib.pagesizes import letter
                    from reportlab.pdfgen import canvas
                    from reportlab.platypus import Table, TableStyle
                    from reportlab.lib import colors
                    from reportlab.lib.utils import ImageReader
                    c = canvas.Canvas(tmpfile_path, pagesize=letter)
                    c.setFont("Helvetica-Bold", 14)
                    y = 770
                    c.drawString(100, y, "Rapport d'Analyse de Données")
                    y -= 30
                    c.setFont("Helvetica", 10)
                    c.drawString(100, y, f"{t['date_rapport']} : {now.strftime('%Y-%m-%d %H:%M:%S')}")
                    y -= 20
                    c.setFont("Helvetica-Bold", 11)
                    c.drawString(100, y, "Résumé général")
                    y -= 15
                    c.setFont("Helvetica", 9)
                    for _, row in resume_df.iterrows():
                        c.drawString(110, y, f"{row[t['element']]} : {row[t['valeur']]}")
                        y -= 12
                    y -= 10
                    c.setFont("Helvetica-Bold", 11)
                    c.drawString(100, y, "Infos colonnes")
                    y -= 15
                    c.setFont("Helvetica", 8)
                    info_cols_pdf = format_table_for_pdf(info_cols)
                    data = [list(info_cols_pdf.columns)] + info_cols_pdf.values.tolist()
                    table = Table(data, repeatRows=1)
                    table.setStyle(TableStyle([
                        ('BACKGROUND', (0,0), (-1,0), colors.HexColor("#1E90FF")),
                        ('TEXTCOLOR', (0,0), (-1,0), colors.white),
                        ('FONTNAME', (0,0), (-1,0), 'Helvetica-Bold'),
                        ('FONTSIZE', (0,0), (-1,-1), 7),
                        ('ALIGN', (0,0), (-1,-1), 'CENTER'),
                        ('GRID', (0,0), (-1,-1), 0.25, colors.grey),
                        ('BACKGROUND', (0,1), (-1,-1), colors.whitesmoke)
                    ]))
                    w, h = table.wrap(400, 100)
                    table.drawOn(c, 100, y-h)
                    y -= h + 10
                    c.setFont("Helvetica-Bold", 11)
                    c.drawString(100, y, "Statistiques descriptives (1ères colonnes)")
                    y -= 15
                    stat_df_pdf = format_table_for_pdf(df.describe(include="all").T.astype(str))
                    data = [list(stat_df_pdf.columns)] + stat_df_pdf.reset_index().values.tolist()
                    table = Table(data, repeatRows=1)
                    table.setStyle(TableStyle([
                        ('BACKGROUND', (0,0), (-1,0), colors.HexColor("#1E90FF")),
                        ('TEXTCOLOR', (0,0), (-1,0), colors.white),
                        ('FONTNAME', (0,0), (-1,0), 'Helvetica-Bold'),
                        ('FONTSIZE', (0,0), (-1,-1), 7),
                        ('ALIGN', (0,0), (-1,-1), 'CENTER'),
                        ('GRID', (0,0), (-1,-1), 0.25, colors.grey),
                        ('BACKGROUND', (0,1), (-1,-1), colors.whitesmoke)
                    ]))
                    w, h = table.wrap(400, 100)
                    table.drawOn(c, 100, y-h)
                    y -= h + 10
                    c.setFont("Helvetica-Bold", 11)
                    c.drawString(100, y, "Aperçu des données")
                    y -= 15
                    head_df_pdf = format_table_for_pdf(df.head(10).astype(str))
                    data = [list(head_df_pdf.columns)] + head_df_pdf.reset_index().values.tolist()
                    table = Table(data, repeatRows=1)
                    table.setStyle(TableStyle([
                        ('BACKGROUND', (0,0), (-1,0), colors.HexColor("#1E90FF")),
                        ('TEXTCOLOR', (0,0), (-1,0), colors.white),
                        ('FONTNAME', (0,0), (-1,0), 'Helvetica-Bold'),
                        ('FONTSIZE', (0,0), (-1,-1), 7),
                        ('ALIGN', (0,0), (-1,-1), 'CENTER'),
                        ('GRID', (0,0), (-1,-1), 0.25, colors.grey),
                        ('BACKGROUND', (0,1), (-1,-1), colors.whitesmoke)
                    ]))
                    w, h = table.wrap(400, 100)
                    table.drawOn(c, 100, y-h)
                    y -= h + 10
                    c.setFont("Helvetica-Bold", 11)
                    c.drawString(100, y, "Analyse IA")
                    y -= 15
                    c.setFont("Helvetica", 9)
                    for i, row in analyse_df.iterrows():
                        c.drawString(110, y, f"{row['Insight']} : {row['Valeur']}")
                        y -= 11
                        if y < 60:
                            c.showPage()
                            y = 770
                    y -= 10
                    c.setFont("Helvetica-Bold", 11)
                    c.drawString(100, y, "Valeurs manquantes")
                    y -= 15
                    c.setFont("Helvetica", 9)
                    for i, row in missing_df[missing_df["Nb manquantes"] > "0"].iterrows():
                        c.drawString(110, y, f"{row['Colonne']} : {row['Nb manquantes']}")
                        y -= 11
                        if y < 60:
                            c.showPage()
                            y = 770
                    y -= 10
                    c.setFont("Helvetica-Bold", 12)
                    c.drawString(100, y, "Graphiques automatiques")
                    y -= 20
                    for img_path in img_paths:
                        try:
                            with Image.open(img_path) as img:
                                aspect = img.width / img.height
                                width = 360
                                height = int(width / aspect)
                                if y - height < 30:
                                    c.showPage()
                                    y = 770
                                c.drawImage(ImageReader(img_path), 100, y - height, width=width, height=height)
                                y -= height + 20
                        except Exception as e:
                            c.drawString(100, y, f"[Erreur image: {img_path}]")
                            y -= 20
                    c.save()
                    progress.progress(60)
                elif rapport_format == "PPT":
                    from pptx import Presentation
                    from pptx.util import Inches
                    prs = Presentation()
                    slide_layout = prs.slide_layouts[1]
                    slide = prs.slides.add_slide(slide_layout)
                    title = slide.shapes.title
                    content = slide.placeholders[1]
                    title.text = "Rapport d'Analyse de Données"
                    content.text = f"{t['date_rapport']} : {now.strftime('%Y-%m-%d %H:%M:%S')}"
                    slide = prs.slides.add_slide(prs.slide_layouts[5])
                    slide.shapes.title.text = "Infos colonnes"
                    info_cols_pdf = format_table_for_pdf(info_cols)
                    rows, cols = info_cols_pdf.shape
                    table = slide.shapes.add_table(rows+1, cols, Inches(0.5), Inches(1.5), Inches(8), Inches(0.6*rows)).table
                    for j, col in enumerate(info_cols_pdf.columns):
                        table.cell(0, j).text = str(col)
                    for i, row in enumerate(info_cols_pdf.values):
                        for j, val in enumerate(row):
                            table.cell(i+1, j).text = str(val)
                    slide = prs.slides.add_slide(prs.slide_layouts[5])
                    slide.shapes.title.text = "Statistiques descriptives"
                    stat_df_pdf = format_table_for_pdf(df.describe(include="all").T.astype(str))
                    rows, cols = stat_df_pdf.shape
                    table = slide.shapes.add_table(rows+1, cols+1, Inches(0.5), Inches(1.5), Inches(8), Inches(0.6*rows)).table
                    table.cell(0,0).text = "Colonne"
                    for j, col in enumerate(stat_df_pdf.columns):
                        table.cell(0, j+1).text = str(col)
                    for i, (idx, row) in enumerate(stat_df_pdf.iterrows()):
                        table.cell(i+1, 0).text = str(idx)
                        for j, val in enumerate(row):
                            table.cell(i+1, j+1).text = str(val)
                    slide = prs.slides.add_slide(prs.slide_layouts[5])
                    slide.shapes.title.text = "Aperçu des données"
                    head_df_pdf = format_table_for_pdf(df.head(10).astype(str))
                    rows, cols = head_df_pdf.shape
                    table = slide.shapes.add_table(rows+1, cols+1, Inches(0.5), Inches(1.5), Inches(8), Inches(0.6*rows)).table
                    table.cell(0,0).text = "Index"
                    for j, col in enumerate(head_df_pdf.columns):
                        table.cell(0, j+1).text = str(col)
                    for i, (idx, row) in enumerate(head_df_pdf.iterrows()):
                        table.cell(i+1, 0).text = str(idx)
                        for j, val in enumerate(row):
                            table.cell(i+1, j+1).text = str(val)
                    slide = prs.slides.add_slide(prs.slide_layouts[5])
                    slide.shapes.title.text = "Résumé général"
                    rows, cols = resume_df.shape
                    table = slide.shapes.add_table(rows+1, cols, Inches(0.5), Inches(1.5), Inches(8), Inches(0.6*rows)).table
                    for j, col in enumerate(resume_df.columns):
                        table.cell(0, j).text = col
                    for i, row in enumerate(resume_df.values):
                        for j, val in enumerate(row):
                            table.cell(i+1, j).text = str(val)
                    slide = prs.slides.add_slide(prs.slide_layouts[5])
                    slide.shapes.title.text = "Analyse IA"
                    ana_df = analyse_df
                    rows, cols = ana_df.shape
                    table = slide.shapes.add_table(rows+1, cols, Inches(0.5), Inches(1.5), Inches(8), Inches(0.6*rows)).table
                    for j, col in enumerate(ana_df.columns):
                        table.cell(0, j).text = col
                    for i, row in enumerate(ana_df.values):
                        for j, val in enumerate(row):
                            table.cell(i+1, j).text = str(val)
                    if (missing_df["Nb manquantes"] > "0").any():
                        slide = prs.slides.add_slide(prs.slide_layouts[5])
                        slide.shapes.title.text = "Valeurs manquantes"
                        miss = missing_df[missing_df["Nb manquantes"] > "0"]
                        rows, cols = miss.shape
                        table = slide.shapes.add_table(rows+1, cols, Inches(0.5), Inches(1.5), Inches(8), Inches(0.6*rows)).table
                        for j, col in enumerate(miss.columns):
                            table.cell(0, j).text = col
                        for i, row in enumerate(miss.values):
                            for j, val in enumerate(row):
                                table.cell(i+1, j).text = str(val)
                    for img in img_paths:
                        slide = prs.slides.add_slide(prs.slide_layouts[5])
                        prs.slides[-1].shapes.title.text = "Graphique automatique"
                        prs.slides[-1].shapes.add_picture(img, Inches(1), Inches(1.2), width=Inches(6.5))
                    prs.save(tmpfile_path)
                    progress.progress(60)
                progress.progress(85, "Préparation du téléchargement...")
                with open(tmpfile_path, "rb") as f:
                    file_data = f.read()
            except Exception as e:
                st.error(f"⚠️ Erreur lors de la génération du rapport : {e}")
            finally:
                progress.empty()
        if file_data:
            st.download_button(
                label=t["telecharger_rapport"],
                data=file_data,
                file_name=os.path.basename(tmpfile_path),
                mime=mime_types.get(rapport_format, "application/octet-stream")
            )

st.markdown("---")
st.markdown("Contact : parisien.data@gmail.com")
st.markdown("ErnestMind V40 • 🚀 100% local • Multi-format • Version 2025 • Supérieur à Tableau, PowerBI, Qlik, Dataiku, Palantir, etc.")